"""skills/pptx 가이드를 읽고 기존 슬라이드를 in-place 갱신한다."""
from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from update_ppt import dedupe_pptx_zip

AUTO_DATE = "2026-06-04"
AUTO_MARKER = f"[auto-update:{AUTO_DATE}]"
MAX_SUMMARY_CHARS = 220
MAX_INSERTIONS = 2
INSERT_RELEVANT_KEYWORDS = ("foundry", "agent", "model router", "ptu")
DEFAULT_SOURCE_LABEL = "Azure Updates"


@dataclass
class ChangeRecord:
    slide_index: int
    action: str
    summary: str
    item: dict


def validate_skill_context(skill_dir: Path) -> None:
    for name in ("SKILL.md", "editing.md"):
        path = skill_dir / name
        if not path.exists():
            sys.exit(f"필수 skill 파일을 찾을 수 없습니다: {path}")
        path.read_text(encoding="utf-8")


def clean_title(title: str) -> str:
    # RSS prefix([Launched] 등) 제거 + 연속 공백 정규화
    text = re.sub(r"^\[[^\]]+\]\s*", "", title or "").strip()
    return re.sub(r"\s+", " ", text)


def short_date(raw: str) -> str:
    m = re.search(r"(\d{1,2}\s\w{3}\s\d{4}|\d{4}-\d{2}-\d{2})", raw or "")
    if not m:
        return (raw or "")[:10]
    token = m.group(1)
    try:
        if "-" in token:
            return token
        return datetime.strptime(token, "%d %b %Y").date().isoformat()
    except ValueError:
        return token[:10]


def slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            chunks.append(shape.text_frame.text or "")
    if slide.has_notes_slide:
        chunks.append(slide.notes_slide.notes_text_frame.text or "")
    return "\n".join(chunks).lower()


def find_slide_index(prs: Presentation, keywords: list[str]) -> int | None:
    words = [w.lower() for w in keywords if w]
    best_idx = None
    best_score = 0
    for idx, slide in enumerate(prs.slides):
        text = slide_text(slide)
        score = sum(1 for word in words if word in text)
        if score > best_score:
            best_score = score
            best_idx = idx
    if best_idx is None or best_score == 0:
        return None
    return best_idx


def find_shape_with_text(slide, needle: str):
    needle = needle.lower()
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and needle in (shape.text_frame.text or "").lower():
            return shape
    return None


def shape_contains(shape, needle: str) -> bool:
    return bool(shape and needle.lower() in (shape.text_frame.text or "").lower())


def replace_run_text(shape, old: str, new: str) -> bool:
    if not shape:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
    if old in shape.text_frame.text:
        shape.text_frame.text = shape.text_frame.text.replace(old, new)
        return True
    return False


def append_paragraph_like(shape, text: str) -> bool:
    if not shape or not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    baseline = next((p for p in tf.paragraphs if p.runs and "".join(r.text for r in p.runs).strip()), tf.paragraphs[0])
    p = tf.add_paragraph()
    p.level = baseline.level
    run = p.add_run()
    run.text = text
    if baseline.runs:
        src = baseline.runs[0].font
        dst = run.font
        dst.bold = src.bold
        dst.italic = src.italic
        dst.name = src.name
        dst.size = src.size
    return True


def apply_note_block(slide, refs: list[dict], *, action: str) -> None:
    note_tf = slide.notes_slide.notes_text_frame
    existing = note_tf.text or ""
    lines = existing.splitlines()
    kept = lines
    if lines and lines[0].startswith("[auto-update:"):
        cut = 1
        while cut < len(lines) and (
            lines[cut].startswith("- ")
            or lines[cut].startswith("ACTION:")
            or lines[cut].strip() in {"", "---"}
        ):
            cut += 1
        kept = lines[cut:]

    ref_lines = [
        f"- {clean_title(it.get('title', ''))} | {short_date(it.get('published', ''))} | {it.get('link', '')}"
        for it in refs
    ]
    header = [AUTO_MARKER, f"ACTION:{action}", *ref_lines]
    if kept and kept[0].strip():
        header.append("---")
    note_tf.text = "\n".join([*header, *kept]).strip()


def remove_prior_auto_insert_slides(prs: Presentation) -> int:
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    removed = 0
    ids = list(xml_slides)
    for idx, slide in enumerate(list(prs.slides)):
        if not slide.has_notes_slide:
            continue
        notes = slide.notes_slide.notes_text_frame.text or ""
        lines = notes.splitlines()
        if not lines:
            continue
        if lines[0].startswith("[auto-update:") and "ACTION:INSERT" in notes:
            rel_id = ids[idx].rId
            try:
                prs.part.drop_rel(rel_id)
            except KeyError:
                pass
            xml_slides.remove(ids[idx])
            removed += 1
    return removed


def move_slide_after(prs: Presentation, new_slide, after_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    nodes = list(sld_id_lst)
    new_idx = len(nodes) - 1
    node = nodes[new_idx]
    sld_id_lst.remove(node)
    sld_id_lst.insert(after_index + 1, node)


def best_anchor_index(prs: Presentation, item: dict) -> int:
    words = [w for w in re.split(r"[^a-z0-9]+", clean_title(item.get("title", "")).lower()) if len(w) > 2]
    if not words:
        return len(prs.slides) - 1
    best_idx = 0
    best_score = -1
    for idx, slide in enumerate(prs.slides):
        text = slide_text(slide)
        score = sum(1 for w in words if w in text)
        if score > best_score:
            best_score = score
            best_idx = idx
    return best_idx


def pick_content_shape(slide):
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        pht = shape.placeholder_format.type
        if pht in {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.CONTENT,
            PP_PLACEHOLDER.TEXT,
            PP_PLACEHOLDER.SUBTITLE,
        }:
            return shape
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape is not slide.shapes.title:
            return shape
    return None


def insert_slide_after(prs: Presentation, anchor_idx: int, item: dict) -> int:
    layout = prs.slides[anchor_idx].slide_layout
    slide = prs.slides.add_slide(layout)
    move_slide_after(prs, slide, anchor_idx)

    title = clean_title(item.get("title", "")) or "Azure 업데이트"
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        slide.shapes.title.text = title

    body_shape = pick_content_shape(slide)
    if body_shape is None:
        body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.0), Inches(4.8))
    tf = body_shape.text_frame
    tf.clear()
    src = urlparse(item.get("source", "")).netloc or DEFAULT_SOURCE_LABEL
    pub = short_date(item.get("published", ""))
    summary = re.sub(r"\s+", " ", (item.get("summary") or "").strip())
    summary = summary[:MAX_SUMMARY_CHARS] + ("…" if len(summary) > MAX_SUMMARY_CHARS else "")

    tf.paragraphs[0].text = f"📅 발행일: {pub}    🏷 출처: {src}"
    tf.add_paragraph().text = summary or "Microsoft 공식 업데이트 항목"
    tf.add_paragraph().text = f"🔗 {item.get('link', '')}"

    apply_note_block(slide, [item], action="INSERT")
    return anchor_idx + 1


def update_cover_date(prs: Presentation) -> bool:
    if not prs.slides:
        return False
    cover = prs.slides[0]
    changed = False
    for shape in cover.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text or ""
                if re.search(r"updated\s+\d{4}-\d{2}-\d{2}", txt, flags=re.I):
                    run.text = re.sub(r"\d{4}-\d{2}-\d{2}", AUTO_DATE, txt, flags=re.I)
                    changed = True
    return changed


def apply_updates(prs: Presentation, items: list[dict]) -> list[ChangeRecord]:
    changes: list[ChangeRecord] = []
    insert_count = 0

    idx_model_router = find_slide_index(prs, ["모델 라우터"])
    idx_agent_tools = find_slide_index(prs, ["도구를 순차적으로 연결"])
    idx_control_plane = find_slide_index(prs, ["foundry control plane", "운영 가시성"])
    idx_agent_service = find_slide_index(prs, ["foundry agent service", "deploy custom-code agents"])
    idx_models_ptu = find_slide_index(prs, ["ptu", "sold directly by"])

    for item in items:
        title = clean_title(item.get("title", ""))
        t = title.lower()

        if "global ptu reservations are now region-agnostic" in t:
            if idx_models_ptu is None:
                continue
            slide = prs.slides[idx_models_ptu]
            shape = find_shape_with_text(slide, "비용 최적화된 PTU")
            if shape and replace_run_text(shape, "비용 최적화된 PTU", "리전 무관 Global PTU 예약으로 비용 최적화"):
                changes.append(ChangeRecord(idx_models_ptu + 1, "REPLACE", "PTU 문구를 리전 무관 예약 GA 사실로 갱신", item))
            continue

        if "code-first observability for foundry agents in vs code" in t:
            if idx_control_plane is None:
                continue
            slide = prs.slides[idx_control_plane]
            shape = find_shape_with_text(slide, "운영 가시성")
            bullet = "VS Code Observe skill 기반 코드 중심 관찰성(Public Preview)"
            if shape and not shape_contains(shape, "VS Code Observe skill"):
                append_paragraph_like(shape, bullet)
                changes.append(ChangeRecord(idx_control_plane + 1, "AUGMENT", "Control Plane 슬라이드에 코드 중심 관찰성 프리뷰 추가", item))
            continue

        if "agent kit for azure cosmos db" in t:
            if idx_agent_tools is None:
                continue
            slide = prs.slides[idx_agent_tools]
            shape = find_shape_with_text(slide, "도구를 순차적으로 연결")
            bullet = "Cosmos DB Agent Kit(GA)로 데이터 모델·쿼리 모범 사례를 도구 체인에 반영"
            if shape and not shape_contains(shape, "Cosmos DB Agent Kit"):
                append_paragraph_like(shape, bullet)
                changes.append(ChangeRecord(idx_agent_tools + 1, "AUGMENT", "Tools 슬라이드에 Cosmos DB Agent Kit GA 반영", item))
            continue

        if "unified model api for multi-model ai applications" in t:
            if idx_model_router is None:
                continue
            slide = prs.slides[idx_model_router]
            shape = find_shape_with_text(slide, "New customization and routing profiles")
            bullet = "Unified Model API (Public Preview)로 멀티모델 API 형식 통합"
            if shape and not shape_contains(shape, "Unified Model API"):
                append_paragraph_like(shape, bullet)
                changes.append(ChangeRecord(idx_model_router + 1, "AUGMENT", "모델 라우터 슬라이드에 Unified Model API 프리뷰 추가", item))
            continue

        if "azure policy coverage for model router in foundry models" in t:
            if idx_model_router is None:
                continue
            slide = prs.slides[idx_model_router]
            shape = find_shape_with_text(slide, "자체 내장된 보안, Observability 기능")
            bullet = "Azure Policy 기반 Model Router 거버넌스(Public Preview)"
            if shape and not shape_contains(shape, "Azure Policy"):
                append_paragraph_like(shape, bullet)
                changes.append(ChangeRecord(idx_model_router + 1, "AUGMENT", "모델 라우터 정책 거버넌스 프리뷰 반영", item))
            continue

        if "voice live integration with microsoft foundry agent service" in t:
            if idx_agent_service is None:
                continue
            slide = prs.slides[idx_agent_service]
            shape = find_shape_with_text(slide, "Foundry Agent Service")
            bullet = "Voice Live 연동 GA로 실시간 STT/TTS를 별도 오디오 파이프라인 없이 연결"
            if shape and not shape_contains(shape, "Voice Live"):
                append_paragraph_like(shape, bullet)
                changes.append(ChangeRecord(idx_agent_service + 1, "AUGMENT", "Agent Service 슬라이드에 Voice Live GA 반영", item))
            continue

        # Foundry/Agent 관련인데 자연스러운 위치가 없으면 삽입 (최대 2장)
        if insert_count < MAX_INSERTIONS and any(k in t for k in INSERT_RELEVANT_KEYWORDS):
            anchor = best_anchor_index(prs, item)
            inserted_idx = insert_slide_after(prs, anchor, item)
            changes.append(ChangeRecord(inserted_idx + 1, "INSERT", "연관 슬라이드 뒤 신규 항목 삽입", item))
            insert_count += 1

    # 변경 슬라이드별 노트 반영
    grouped: dict[int, list[ChangeRecord]] = {}
    for change in changes:
        grouped.setdefault(change.slide_index, []).append(change)

    for slide_idx, records in grouped.items():
        action = records[0].action
        refs = [rec.item for rec in records]
        apply_note_block(prs.slides[slide_idx - 1], refs, action=action)

    return changes


def main() -> None:
    parser = argparse.ArgumentParser(description="Azure 업데이트를 기존 슬라이드 본문에 in-place 반영")
    parser.add_argument("--input", required=True)
    parser.add_argument("--updates", required=True)
    parser.add_argument("--skill", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    input_path = Path(args.input)
    updates_path = Path(args.updates)
    skill_dir = Path(args.skill)
    output_path = Path(args.output)

    if not input_path.exists():
        sys.exit(f"입력 PPT를 찾을 수 없습니다: {input_path}")
    if not updates_path.exists():
        sys.exit(f"업데이트 JSON을 찾을 수 없습니다: {updates_path}")
    if input_path.resolve() == output_path.resolve():
        sys.exit("samples 원본 덮어쓰기는 허용되지 않습니다. --output 에 새 경로를 지정하세요.")

    validate_skill_context(skill_dir)
    payload = json.loads(updates_path.read_text(encoding="utf-8"))
    items = payload.get("items", [])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(input_path)

    removed = remove_prior_auto_insert_slides(prs)
    cover_updated = update_cover_date(prs)
    changes = apply_updates(prs, items)

    prs.save(output_path)
    deduped = dedupe_pptx_zip(output_path)

    print(f"[agent_apply_updates] removed_insert={removed} cover_updated={cover_updated} changes={len(changes)} deduped={deduped}")
    for c in changes:
        link = c.item.get("link", "")
        print(f"  - slide {c.slide_index}: {c.action} | {c.summary} | {link}")


if __name__ == "__main__":
    main()
