"""PPTX 내용 -> 스크롤형 HTML 덱 빌드.

릴리즈 폴더(releases/YYYYMMDD-HHMMSS-<자료명>/) 의 .pptx 에서 슬라이드의
**텍스트 내용** 을 추출해 슬라이드 이미지가 아닌 네이티브 HTML 로 렌더링한다.

레이아웃:
  - 왼쪽 사이드바: 슬라이드 목차(타이틀/번호) + scroll-spy 하이라이트 + 클릭 스크롤
  - 본문: 위 -> 아래로 세로 스크롤되는 슬라이드 섹션

타이틀 우선순위:
  1. 릴리즈 폴더의 titles.json  {"1": "제목", "2": "제목", ...}  (AI 가 슬라이드 내용 보고 작성)
  2. PPTX 에서 추출한 슬라이드 제목/첫 텍스트
  3. "Slide N"

버전 히스토리:
  릴리즈 타임스탬프 = 한 버전. 자료명(stem) 이 같은 릴리즈들이 한 덱의 버전 묶음.
  - docs/decks/<release-id>.html : 버전별 자체완결 HTML (버전 전환기 포함)
  - docs/decks/index.html        : 전체 덱/버전 히스토리 목록 (최신순)

순수 Python + python-pptx 만 사용한다 (새 의존성 없음).
"""
from __future__ import annotations

import argparse
import base64
import json
import re
import shutil
import zipfile
from collections import Counter
from html import escape
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

RELEASE_NAME_RE = re.compile(r"^(\d{8})-(\d{6})-(.+)$")


# --------------------------------------------------------------------------- #
# 릴리즈 스캔
# --------------------------------------------------------------------------- #
def list_releases(releases_dir: Path) -> list[Path]:
    if not releases_dir.exists():
        return []
    items = [p for p in releases_dir.iterdir() if p.is_dir() and RELEASE_NAME_RE.match(p.name)]
    return sorted(items, key=lambda p: p.name, reverse=True)


def parse_release(folder: Path) -> dict:
    m = RELEASE_NAME_RE.match(folder.name)
    date_s, time_s, stem = m.group(1), m.group(2), m.group(3)
    ts_human = f"{date_s[:4]}-{date_s[4:6]}-{date_s[6:8]} {time_s[:2]}:{time_s[2:4]}:{time_s[4:6]} KST"
    pptx = next((p for p in folder.iterdir() if p.suffix.lower() == ".pptx"), None)
    titles_file = folder / "titles.json"
    content_file = folder / "content.json"
    slides_dir = folder / "slides"
    return {
        "folder": folder,
        "id": folder.name,
        "stem": stem,
        "ts_human": ts_human,
        "ts_key": f"{date_s}-{time_s}",
        "pptx": pptx,
        "titles_file": titles_file if titles_file.exists() else None,
        "content_file": content_file if content_file.exists() else None,
        "slides_dir": slides_dir if slides_dir.is_dir() else None,
    }


# --------------------------------------------------------------------------- #
# PPTX 구조 추출 (위치·크기·폰트·색상·이미지·표)
# --------------------------------------------------------------------------- #
EMU_PER_PT = 12700


def _clean(text: str) -> str:
    return text.replace("\x0b", "\n").replace("\r", "\n")


def _align(a) -> str:
    return {1: "left", 2: "center", 3: "right", 4: "justify"}.get(
        int(a) if a is not None else 0, "left")


def _valign(tf) -> str:
    try:
        v = tf.vertical_anchor
    except Exception:
        v = None
    return {3: "center", 4: "flex-end"}.get(int(v) if v is not None else 1, "flex-start")


def _solid_rgb(color_fmt) -> str | None:
    try:
        if color_fmt.type is not None and color_fmt.rgb is not None:
            return "#" + str(color_fmt.rgb)
    except Exception:
        pass
    return None


def _fill_rgb(sh) -> str | None:
    try:
        if int(sh.fill.type) == 1:  # solid
            return _solid_rgb(sh.fill.fore_color)
    except Exception:
        pass
    return None


def _line_rgb(sh) -> str | None:
    try:
        return _solid_rgb(sh.line.color)
    except Exception:
        return None


def _runs(paragraph, default_pt: float) -> list[dict]:
    runs = []
    src = paragraph.runs or []
    for r in src:
        size = r.font.size
        pt = (size / EMU_PER_PT) if size is not None else default_pt
        runs.append({
            "text": _clean(r.text),
            "pt": pt,
            "bold": bool(r.font.bold),
            "italic": bool(r.font.italic),
            "color": _solid_rgb(r.font.color),
        })
    if not src and paragraph.text.strip():
        runs.append({"text": _clean(paragraph.text), "pt": default_pt,
                     "bold": False, "italic": False, "color": None})
    return runs


def _text_shape(sh, default_pt: float) -> dict | None:
    tf = sh.text_frame
    paras = []
    for p in tf.paragraphs:
        runs = _runs(p, default_pt)
        if any(r["text"].strip() for r in runs):
            paras.append({"align": _align(p.alignment), "runs": runs})
    if not paras:
        return None
    return {"kind": "text", "valign": _valign(tf), "bg": _fill_rgb(sh), "paras": paras}


def _image_data_uri(sh) -> str | None:
    try:
        img = sh.image
        b64 = base64.b64encode(img.blob).decode("ascii")
        return f"data:{img.content_type};base64,{b64}"
    except Exception:
        return None


def _table_shape(sh) -> dict | None:
    rows = []
    for r in sh.table.rows:
        rows.append([_clean(c.text).strip() for c in r.cells])
    if not any(any(c for c in row) for row in rows):
        return None
    return {"kind": "table", "rows": rows}


def _norm(s: str) -> str:
    return " ".join((s or "").split()).strip().lower()


def _present_paras(sh) -> list[dict]:
    """프레젠테이션용: 문단 = 한 줄 불릿. 들여쓰기 레벨/볼드/정렬 보존."""
    out = []
    for p in sh.text_frame.paragraphs:
        txt = _clean(p.text).strip()
        if not txt:
            continue
        runs = p.runs or []
        bold = bool(runs) and all(bool(r.font.bold) for r in runs)
        out.append({
            "text": txt,
            "level": int(p.level or 0),
            "bold": bold,
            "align": _align(p.alignment),
        })
    return out


def _video_url_info(url: str) -> dict | None:
    """영상 URL → {kind:'embed'|'file', src, provider}. YouTube/Vimeo 는 iframe 임베드로."""
    if not url:
        return None
    u = url.strip()
    yt = re.search(r"(?:youtube\.com/(?:watch\?v=|embed/|v/)|youtu\.be/)([\w-]{6,})", u)
    if yt:
        vid = yt.group(1)
        start = re.search(r"[?&](?:start|t)=(\d+)", u)
        q = f"?start={start.group(1)}" if start else ""
        return {"kind": "embed", "src": f"https://www.youtube.com/embed/{vid}{q}", "provider": "youtube"}
    vm = re.search(r"vimeo\.com/(?:video/)?(\d+)", u)
    if vm:
        return {"kind": "embed", "src": f"https://player.vimeo.com/video/{vm.group(1)}", "provider": "vimeo"}
    if re.search(r"\.(mp4|webm|ogg|mov|m4v)(\?|$)", u, re.I) or u.startswith("data:video"):
        return {"kind": "file", "src": u, "provider": "file"}
    # 알 수 없는 http(s) 링크는 임베드로 시도
    if u.startswith("http"):
        return {"kind": "embed", "src": u, "provider": "external"}
    return {"kind": "file", "src": u, "provider": "file"}


def _video_html(b: dict) -> str:
    """video/embed 블록 → 반응형 영상 플레이어."""
    info = _video_url_info(b.get("src") or b.get("url") or "")
    if not info:
        return ""
    cap = f"<figcaption>{_etext(b['caption'])}</figcaption>" if b.get("caption") else ""
    if info["kind"] == "embed":
        inner = (f'<iframe src="{escape(info["src"])}" title="video" loading="lazy" '
                 f'allow="accelerometer; autoplay; clipboard-write; encrypted-media; '
                 f'gyroscope; picture-in-picture; web-share" '
                 f'allowfullscreen></iframe>')
    else:
        attrs = " controls"
        if b.get("autoplay"):
            attrs += " autoplay muted"
        if b.get("loop"):
            attrs += " loop"
        if b.get("muted"):
            attrs += " muted"
        poster = f' poster="{escape(b["poster"])}"' if b.get("poster") else ""
        inner = f'<video{attrs}{poster}><source src="{escape(info["src"])}"></video>'
    return f'<figure class="video"><div class="frame">{inner}</div>{cap}</figure>'


def _slide_videos(slide) -> list[dict]:
    """슬라이드의 임베디드/링크 영상을 video 블록 리스트로. (YouTube 등 외부 링크, 내장 파일)"""
    out = []
    try:
        rels = slide.part.rels
    except Exception:
        return out
    for rel in rels.values():
        try:
            if "video" not in rel.reltype:
                continue
            if rel.is_external:
                out.append({"type": "video", "src": rel.target_ref})
            else:
                part = rel.target_part
                b64 = base64.b64encode(part.blob).decode("ascii")
                out.append({"type": "video", "src": f"data:{part.content_type};base64,{b64}"})
        except Exception:
            continue
    return out


def _image_src(sh, max_w: int = 1000) -> str | None:
    """임베디드 이미지를 data URI 로. Pillow 있으면 축소·WebP 로 용량 절감.
    애니메이션 GIF 는 프레임 보존을 위해 변환하지 않고 원본 그대로 인라인."""
    try:
        blob = sh.image.blob
        ctype = sh.image.content_type
    except Exception:
        return None
    if ctype == "image/gif":
        b64 = base64.b64encode(blob).decode("ascii")
        return f"data:image/gif;base64,{b64}"
    try:
        from io import BytesIO
        from PIL import Image
        im = Image.open(BytesIO(blob))
        has_alpha = im.mode in ("RGBA", "LA", "P")
        im = im.convert("RGBA") if has_alpha else im.convert("RGB")
        if im.width > max_w:
            ratio = max_w / im.width
            im = im.resize((max_w, round(im.height * ratio)))
        buf = BytesIO()
        if im.mode == "RGBA":
            im.save(buf, "WEBP", quality=85, method=4)
        else:
            im.save(buf, "WEBP", quality=82, method=4)
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        return f"data:image/webp;base64,{b64}"
    except Exception:
        b64 = base64.b64encode(blob).decode("ascii")
        return f"data:{ctype};base64,{b64}"


def _geom(left, top, width, height, sw, sh_) -> dict:
    return {
        "l": (left or 0) / sw * 100,
        "t": (top or 0) / sh_ * 100,
        "w": (width or 0) / sw * 100,
        "h": (height or 0) / sh_ * 100,
    }


def _walk_shapes(shapes, tf_xform, out):
    """그룹 변환을 누적해 모든 도형을 (절대 EMU 좌표와 함께) 평탄화."""
    ox, oy, sx, sy = tf_xform
    for sh in shapes:
        is_group = sh.shape_type is not None and int(sh.shape_type) == 6
        if is_group:
            gl, gt = sh.left or 0, sh.top or 0
            gw, gh = sh.width or 0, sh.height or 0
            try:
                xf = sh._element.grpSpPr.xfrm
                cox, coy = xf.chOff.x, xf.chOff.y
                cex, cey = xf.chExt.cx, xf.chExt.cy
            except Exception:
                cox = coy = 0
                cex, cey = gw or 1, gh or 1
            groot_l = ox + gl * sx
            groot_t = oy + gt * sy
            nsx = (gw * sx) / (cex or 1)
            nsy = (gh * sy) / (cey or 1)
            _walk_shapes(sh.shapes, (groot_l - cox * nsx, groot_t - coy * nsy, nsx, nsy), out)
            continue
        l = ox + (sh.left or 0) * sx
        t = oy + (sh.top or 0) * sy
        w = (sh.width or 0) * sx
        h = (sh.height or 0) * sy
        out.append((sh, l, t, w, h))


def extract_deck(pptx_path: Path) -> list[dict]:
    """슬라이드별 {num, auto_title, w_pt, h_pt, shapes[]} (구조 보존)."""
    prs = Presentation(str(pptx_path))
    sw, sh_ = prs.slide_width, prs.slide_height
    w_pt = sw / EMU_PER_PT
    deck = []
    for i, slide in enumerate(prs.slides, 1):
        flat: list = []
        _walk_shapes(slide.shapes, (0, 0, 1, 1), flat)

        title_shape = slide.shapes.title
        auto_title = ""
        if title_shape is not None:
            auto_title = " ".join(_clean(title_shape.text).split()).strip()

        shapes_out = []
        for sh, l, t, w, h in flat:
            geom = _geom(l, t, w, h, sw, sh_)
            item = None
            if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                src = _image_data_uri(sh)
                if src:
                    item = {"kind": "image", "src": src}
            elif sh.has_table:
                item = _table_shape(sh)
            elif sh.has_text_frame and sh.text_frame.text.strip():
                is_title = title_shape is not None and sh == title_shape
                default_pt = 28.0 if is_title else 14.0
                item = _text_shape(sh, default_pt)
            else:
                fill = _fill_rgb(sh)
                line = _line_rgb(sh)
                if fill or line:
                    item = {"kind": "box", "bg": fill, "border": line}
            if item is None:
                continue
            item.update(geom)
            shapes_out.append(item)

        if not auto_title:
            for it in shapes_out:
                if it.get("kind") == "text":
                    auto_title = " ".join(
                        r["text"] for p in it["paras"] for r in p["runs"]).split("\n")[0].strip()
                    if auto_title:
                        break

        # 프레젠테이션용 블록: 읽기 순서(위→아래, 좌→우)로 불릿/표/이미지
        title_norm = _norm(auto_title)
        ordered = sorted(flat, key=lambda x: (round(x[2] / sh_, 3), round(x[1] / sw, 3)))
        blocks = []
        for sh, l, t, w, h in ordered:
            if title_shape is not None and sh == title_shape:
                continue
            if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                src = _image_src(sh)
                if src:
                    blocks.append({"type": "image", "src": src})
            elif sh.has_table:
                tb = _table_shape(sh)
                if tb:
                    blocks.append({"type": "table", "rows": tb["rows"]})
            elif sh.has_text_frame and sh.text_frame.text.strip():
                items = _present_paras(sh)
                # 본문에 제목이 그대로 다시 나오면 제거(중복 방지)
                if items and title_norm and _norm(items[0]["text"]) == title_norm:
                    items = items[1:]
                if items:
                    blocks.append({"type": "bullets", "items": items})

        # 임베디드/링크 영상 (YouTube 등) → video 블록
        blocks.extend(_slide_videos(slide))

        # 발표자 노트(슬라이드 노트) → 발표자 모드용 스크립트
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = _clean(slide.notes_slide.notes_text_frame.text).strip()
        except Exception:
            notes = ""

        deck.append({"num": i, "auto_title": auto_title, "w_pt": w_pt,
                     "ar": sw / sh_, "shapes": shapes_out, "blocks": blocks,
                     "notes": notes})
    return deck


def load_titles(titles_file: Path | None) -> dict:
    if not titles_file:
        return {}
    try:
        data = json.loads(titles_file.read_text(encoding="utf-8"))
        return {str(k): str(v) for k, v in data.items()}
    except Exception:
        return {}


def resolve_title(slide: dict, overrides: dict) -> str:
    key = str(slide["num"])
    if overrides.get(key, "").strip():
        return overrides[key].strip()
    if slide["auto_title"].strip():
        return slide["auto_title"].strip()
    return f"Slide {slide['num']}"


# --------------------------------------------------------------------------- #
# 테마 추출 (소스 PPT 의 색/배경에서 덱 테마를 도출)
# --------------------------------------------------------------------------- #
def _hex_lum(h: str) -> float:
    h = h.lstrip("#")
    r, g, b = (int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _hex_sat(h: str) -> float:
    h = h.lstrip("#")
    r, g, b = (int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))
    mx, mn = max(r, g, b), min(r, g, b)
    return 0.0 if mx == 0 else (mx - mn) / mx


def _mix(hex_a: str, hex_b: str, t: float) -> str:
    a = hex_a.lstrip("#"); b = hex_b.lstrip("#")
    out = []
    for i in (0, 2, 4):
        ca, cb = int(a[i:i + 2], 16), int(b[i:i + 2], 16)
        out.append(round(ca + (cb - ca) * t))
    return "#" + "".join(f"{c:02X}" for c in out)


# Foundry(보라) 기본 테마 — 추출 실패 시 폴백
_FALLBACK_THEME = {
    "mode": "dark", "paper": "#0e0b1c", "ink": "#f4f3fb",
    "navy": "#0b0a14", "section": "#140f28",
    "accent": "#8b5cf6", "accent_deep": "#6d28d9", "accent_cyan": "#e3008c",
    "accent_soft": "#a78bfa",
}


def extract_theme(pptx_path: Path) -> dict:
    """PPTX 의 슬라이드 배경/색 사용을 스캔해 덱 테마(밝기·포인트색)를 도출한다.

    - 슬라이드 배경의 밝기 분포로 light/dark 기본 모드 결정
    - 가장 많이 쓰인 채도 높은 색을 포인트(accent)로 채택
    - 어두운 배경색(타이틀/섹션용), 밝은 배경색(본문용) 추출
    """
    try:
        z = zipfile.ZipFile(str(pptx_path))
    except Exception:
        return dict(_FALLBACK_THEME)
    names = z.namelist()

    def rd(n: str) -> str:
        try:
            return z.read(n).decode("utf-8", "ignore")
        except Exception:
            return ""

    theme_acc: list[str] = []
    for n in names:
        if n.startswith("ppt/theme/theme") and n.endswith(".xml"):
            tx = rd(n)
            for tag in ("accent1", "accent2", "accent3", "accent4", "accent5", "accent6"):
                m = re.search(rf'<a:{tag}>\s*<a:srgbClr val="([0-9A-Fa-f]{{6}})"', tx)
                if m:
                    theme_acc.append("#" + m.group(1).upper())
            break

    bg_light: Counter = Counter()
    bg_dark: Counter = Counter()
    used: Counter = Counter()
    slide_names = sorted(n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n))
    for n in slide_names:
        sx = rd(n)
        mbg = re.search(r"<p:bg>.*?</p:bg>", sx, re.S)
        if mbg:
            cols = re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', mbg.group(0))
            if cols:
                c = "#" + cols[0].upper()
                (bg_light if _hex_lum(c) >= 0.5 else bg_dark)[c] += 1
        for c in re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', sx):
            used["#" + c.upper()] += 1

    light_total = sum(bg_light.values())
    dark_total = sum(bg_dark.values())
    if light_total == 0 and dark_total == 0:
        return dict(_FALLBACK_THEME)
    mode = "light" if light_total >= dark_total else "dark"

    paper = bg_light.most_common(1)[0][0] if bg_light else "#F7FAFC"
    darks = sorted(bg_dark, key=_hex_lum)
    navy = darks[0] if darks else "#0B1F3A"
    section = darks[-1] if darks else navy

    def good_accent(c: str) -> bool:
        return _hex_sat(c) >= 0.32 and 0.10 <= _hex_lum(c) <= 0.86

    cand = [c for c, _ in used.most_common() if good_accent(c)]
    for c in theme_acc:
        if good_accent(c) and c not in cand:
            cand.append(c)
    if not cand:
        cand = ["#0078D4", "#50E6FF", "#004578"]

    accent = cand[0]
    # 깊은 포인트색: 채도 있는 어두운 후보 우선, 없으면 accent 를 어둡게
    dark_cands = [c for c in cand if 0.10 <= _hex_lum(c) <= 0.34]
    accent_deep = dark_cands[0] if dark_cands else _mix(accent, "#001022", 0.42)
    # 밝은 보조 포인트색(시안 계열): 채도 있는 밝은 후보 우선, 없으면 accent 를 밝게
    bright_cands = [c for c in cand if 0.58 <= _hex_lum(c) <= 0.92]
    accent_cyan = bright_cands[0] if bright_cands else _mix(accent, "#FFFFFF", 0.5)
    accent_soft = _mix(accent, "#FFFFFF", 0.4)

    if not darks:
        navy = _mix(accent_deep, "#000814", 0.55)
        section = _mix(accent_deep, "#0A1426", 0.4)

    return {
        "mode": mode,
        "paper": paper,
        "ink": "#0B1F3A" if mode == "light" else "#eef5fc",
        "navy": navy,
        "section": section,
        "accent": accent,
        "accent_deep": accent_deep,
        "accent_cyan": accent_cyan,
        "accent_soft": accent_soft,
    }


def render_theme_css(t: dict) -> str:
    """추출 테마를 .frag / 셸 CSS 토큰 오버라이드로 변환."""
    a = t["accent"]; deep = t["accent_deep"]; cyan = t["accent_cyan"]; soft = t["accent_soft"]
    navy = t["navy"]; section = t["section"]; paper = t["paper"]; ink = t["ink"]
    light = t["mode"] == "light"
    return f"""
/* === 소스 PPT 에서 추출한 덱 테마 === */
:root{{--accent:{a};--active:{_mix(a, '#FFFFFF', 0.84)};}}
.slide-wrap .frag{{--msft-blue:{a};--msft-purple:{deep};--msft-pink:{a};}}
.frag--dark{{--msft-blue:{cyan};--msft-purple:{soft};--msft-pink:{cyan};
background:{navy};color:#eef5fc;background-image:
 radial-gradient(120% 90% at 100% 0%,color-mix(in srgb,{a} 26%,transparent),transparent 55%),
 radial-gradient(95% 80% at 0% 100%,color-mix(in srgb,{cyan} 16%,transparent),transparent 60%),
 linear-gradient(180deg,{navy},{section});}}
.frag--light{{--msft-blue:{a};--msft-purple:{deep};--msft-pink:{deep};
background:{paper};color:{ink if light else '#0B1F3A'};background-image:
 radial-gradient(120% 90% at 100% 0%,color-mix(in srgb,{a} 12%,transparent),transparent 58%),
 radial-gradient(90% 80% at 0% 100%,color-mix(in srgb,{cyan} 9%,transparent),transparent 62%);}}
.frag--plain{{--msft-blue:{a};--msft-purple:{deep};--msft-pink:{deep};}}
/* 다이어그램은 항상 짙은 네이비 카드 위에 (밝은 슬라이드에서도 또렷하게) */
.frag .fdiagram .mermaid{{background:{navy};border:1px solid color-mix(in srgb,{a} 30%,transparent);
border-radius:14px;padding:clamp(14px,2vw,22px);box-shadow:0 10px 30px rgba(7,16,30,.28);}}
"""


def mermaid_theme_json(t: dict) -> str:
    """Mermaid 초기화용 themeVariables (네이비 카드 + 추출 포인트색)."""
    a = t["accent"]; cyan = t["accent_cyan"]; navy = t["navy"]
    return json.dumps({
        "background": "transparent",
        "primaryColor": _mix(navy, "#FFFFFF", 0.10),
        "primaryBorderColor": a,
        "primaryTextColor": "#eaf2fb",
        "secondaryColor": _mix(navy, a, 0.22),
        "tertiaryColor": _mix(navy, "#FFFFFF", 0.04),
        "lineColor": cyan,
        "edgeLabelBackground": _mix(navy, "#000000", 0.25),
        "clusterBkg": _mix(navy, a, 0.16),
        "clusterBorder": a,
        "fontSize": "15px",
    }, ensure_ascii=False)


# --------------------------------------------------------------------------- #
# HTML 렌더
# --------------------------------------------------------------------------- #
CSS = """:root{color-scheme:light dark;--bg:#fff;--fg:#0f172a;--muted:#64748b;
--border:#e2e8f0;--card:#fff;--accent:#2563eb;--side:#f8fafc;--active:#e0e7ff;
--slide-bg:#ffffff;--stage-shadow:0 1px 3px rgba(0,0,0,.12),0 8px 24px rgba(0,0,0,.08);}
@media (prefers-color-scheme:dark){:root{--bg:#0b1120;--fg:#e2e8f0;--muted:#94a3b8;
--border:#1e293b;--card:#0f172a;--accent:#60a5fa;--side:#0f172a;--active:#1e293b;
--slide-bg:#ffffff;--stage-shadow:0 1px 3px rgba(0,0,0,.5),0 8px 24px rgba(0,0,0,.4);}}
*{box-sizing:border-box;}
html{scroll-behavior:smooth;}
body{margin:0;font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,
"Helvetica Neue",Arial,"Apple SD Gothic Neo","Noto Sans KR",sans-serif;
background:var(--bg);color:var(--fg);line-height:1.4;}
.layout{display:grid;grid-template-columns:280px 1fr;min-height:100vh;transition:grid-template-columns .25s ease;}
body.side-collapsed .layout{grid-template-columns:0 1fr;}
aside{position:sticky;top:0;align-self:start;height:100vh;overflow-y:auto;
background:var(--side);border-right:1px solid var(--border);padding:20px 0;}
body.side-collapsed aside{overflow:hidden;border-right:none;}
.collapse-btn{position:absolute;top:14px;right:12px;width:30px;height:30px;display:grid;
place-items:center;border:1px solid var(--border);border-radius:8px;background:var(--card);
color:var(--fg);cursor:pointer;font-size:14px;line-height:1;z-index:5;}
.collapse-btn:hover{background:var(--active);}
.side-open{position:fixed;top:14px;left:14px;z-index:30;width:38px;height:38px;display:none;
place-items:center;border:1px solid var(--border);border-radius:10px;background:var(--card);
color:var(--fg);cursor:pointer;font-size:18px;box-shadow:0 4px 16px rgba(0,0,0,.18);}
body.side-collapsed .side-open{display:grid;}
aside .deck-title{padding:0 50px 12px 20px;font-weight:700;font-size:15px;}
aside .ver{padding:0 20px 16px;font-size:12px;color:var(--muted);}
aside .ver select{width:100%;margin-top:6px;padding:6px;border-radius:8px;
border:1px solid var(--border);background:var(--card);color:var(--fg);font-size:12px;}
nav ol{list-style:none;margin:0;padding:0;}
nav a{display:flex;gap:10px;padding:8px 20px;text-decoration:none;color:var(--fg);
font-size:13px;border-left:3px solid transparent;}
nav a:hover{background:var(--active);}
nav a.active{background:var(--active);border-left-color:var(--accent);font-weight:600;}
nav a .n{color:var(--muted);min-width:22px;text-align:right;font-variant-numeric:tabular-nums;}
main{padding:32px clamp(16px,4vw,64px) 40vh;}
.slide-wrap{max-width:1180px;margin:0 auto 26px;scroll-margin-top:20px;
position:relative;aspect-ratio:16/9;}
.card{background:var(--card);border:1px solid var(--border);border-radius:16px;
box-shadow:var(--stage-shadow);padding:clamp(24px,3.5vw,48px);min-height:60vh;
display:flex;flex-direction:column;}
.card-head{margin:0 0 22px;padding-bottom:16px;border-bottom:2px solid var(--accent);}
.card-head .kicker{color:var(--accent);font-weight:700;font-size:12px;
letter-spacing:.06em;text-transform:uppercase;font-variant-numeric:tabular-nums;}
.card-head h2{margin:6px 0 0;font-size:clamp(22px,2.8vw,34px);line-height:1.22;
font-weight:750;}
.card-body{flex:1;font-size:clamp(15px,1.45vw,19px);}
.card-body>*:first-child{margin-top:0;}
ul.bullets{list-style:none;margin:14px 0;padding:0;}
ul.bullets li{position:relative;margin:9px 0;padding-left:1.4em;line-height:1.5;}
ul.bullets li::before{content:"\\2022";position:absolute;left:.25em;color:var(--accent);
font-weight:700;}
ul.bullets li.lv1{margin-left:1.5em;font-size:.96em;}
ul.bullets li.lv2{margin-left:3em;font-size:.93em;}
ul.bullets li.lv3,ul.bullets li.lv4{margin-left:4.5em;font-size:.9em;opacity:.95;}
ul.bullets li.lv1::before{content:"\\2013";}
ul.bullets li.lv2::before,ul.bullets li.lv3::before,ul.bullets li.lv4::before{content:"\\25E6";}
ul.bullets li.b{font-weight:700;}
.card-body .tbl{overflow-x:auto;margin:18px 0;}
.card-body table{border-collapse:collapse;width:100%;font-size:.92em;}
.card-body th,.card-body td{border:1px solid var(--border);padding:8px 12px;
vertical-align:top;text-align:left;}
.card-body table tr:first-child td{background:var(--side);font-weight:600;}
.card-body figure.img{margin:18px 0;text-align:center;}
.card-body figure.img img{max-width:100%;max-height:52vh;border-radius:8px;
border:1px solid var(--border);}
/* 영상 */
figure.video{margin:18px 0;}
figure.video .frame{position:relative;width:100%;max-width:880px;margin-inline:auto;
aspect-ratio:16/9;background:#000;border-radius:10px;overflow:hidden;border:1px solid var(--border);}
figure.video .frame iframe,figure.video .frame video{position:absolute;inset:0;width:100%;height:100%;border:0;}
figure.video figcaption{margin-top:8px;text-align:center;font-size:.88em;color:var(--muted);}
.empty{color:var(--muted);font-style:italic;}
/* 발표용(에이전트 재구성) 블록 */
.card-body .lead{font-size:1.12em;line-height:1.62;margin:0 0 18px;font-weight:500;}
.card-body p{line-height:1.62;margin:12px 0;}
ul.bul{list-style:none;margin:14px 0;padding:0;}
ul.bul>li{position:relative;margin:9px 0;padding-left:1.4em;line-height:1.55;}
ul.bul>li::before{content:"\\2022";position:absolute;left:.25em;color:var(--accent);font-weight:700;}
ul.bul ul{list-style:none;margin:6px 0 2px;padding:0;}
ul.bul ul li{position:relative;margin:5px 0;padding-left:1.3em;font-size:.94em;color:var(--muted);}
ul.bul ul li::before{content:"\\2013";position:absolute;left:.2em;color:var(--muted);}
.grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(210px,1fr));gap:14px;margin:16px 0;}
.gc{border:1px solid var(--border);border-radius:12px;padding:16px 18px;background:var(--side);}
.gc h4{margin:0 0 6px;font-size:1.02em;color:var(--accent);}
.gc p{margin:0;font-size:.93em;color:var(--muted);line-height:1.5;}
ol.steps{list-style:none;margin:16px 0;padding:0;}
ol.steps>li{display:flex;gap:14px;margin:14px 0;align-items:flex-start;}
ol.steps .sn{flex:none;width:30px;height:30px;border-radius:50%;background:var(--accent);
color:#fff;display:grid;place-items:center;font-weight:700;font-size:.9em;}
ol.steps h4{margin:2px 0 4px;font-size:1.02em;}
ol.steps p{margin:0;color:var(--muted);font-size:.94em;line-height:1.5;}
.flow{display:flex;flex-wrap:wrap;align-items:center;gap:6px;margin:18px 0;}
.flow .fi{border:1px solid var(--border);background:var(--side);border-radius:10px;
padding:10px 14px;font-weight:600;font-size:.94em;}
.flow .fa{color:var(--accent);font-weight:700;font-size:1.1em;}
.stats{display:flex;flex-wrap:wrap;gap:16px;margin:18px 0;}
.stat{background:var(--side);border:1px solid var(--border);border-radius:12px;padding:16px 22px;min-width:120px;}
.stat .sv{font-size:1.9em;font-weight:750;color:var(--accent);line-height:1;}
.stat .sl{margin-top:6px;font-size:.9em;color:var(--muted);}
.callout{border-left:4px solid var(--accent);background:var(--side);padding:12px 16px;
border-radius:0 10px 10px 0;margin:16px 0;line-height:1.55;}
.callout.tip{border-left-color:#16a34a;} .callout.warn{border-left-color:#d97706;}
.cols{display:grid;grid-template-columns:repeat(auto-fit,minmax(240px,1fr));gap:20px;margin:16px 0;}
.col h4{margin:0 0 8px;color:var(--accent);border-bottom:1px solid var(--border);padding-bottom:6px;}
blockquote{margin:16px 0;padding:14px 20px;border-left:4px solid var(--accent);
font-style:italic;background:var(--side);border-radius:0 10px 10px 0;line-height:1.55;}
blockquote cite{display:block;margin-top:8px;font-style:normal;font-size:.88em;color:var(--muted);}
.card-body table th{background:var(--accent);color:#fff;border-color:var(--accent);font-weight:600;}
.card--hero{align-items:center;justify-content:center;text-align:center;min-height:72vh;}
.card--hero .card-head{border-bottom:none;margin:0;padding:0;}
.card--hero .card-head h2{font-size:clamp(28px,4vw,44px);margin-top:12px;line-height:1.15;}
.card--hero .card-body{flex:none;margin-top:18px;}
.card--hero .lead{font-size:1.2em;color:var(--muted);font-weight:400;}
.topbar{display:none;}
@media (max-width:860px){
.layout{grid-template-columns:1fr;}
body.side-collapsed .layout{grid-template-columns:1fr;}
.collapse-btn,.side-open{display:none!important;}
main{padding:16px 12px 40vh;}
aside{position:fixed;left:0;top:0;width:80%;max-width:320px;z-index:20;
transform:translateX(-100%);transition:transform .2s;box-shadow:0 0 40px rgba(0,0,0,.3);}
aside.open{transform:translateX(0);}
.topbar{display:flex;position:sticky;top:0;z-index:10;background:var(--side);
border-bottom:1px solid var(--border);padding:10px 16px;gap:12px;align-items:center;}
.topbar button{font-size:20px;background:none;border:none;color:var(--fg);cursor:pointer;}
.topbar .t{font-weight:700;font-size:14px;}
}
/* === 모델 작성 슬라이드 조각(.frag) 디자인 시스템 === */
.slide-wrap .frag{position:absolute;inset:0;border-radius:16px;overflow:hidden;
box-shadow:var(--stage-shadow);padding:clamp(22px,3.2vw,46px);
display:flex;flex-direction:column;
--msft-purple:#8b5cf6;--msft-blue:#3b82f6;--msft-pink:#e3008c;
--ink-dim:color-mix(in srgb,currentColor 60%,transparent);
--hair:color-mix(in srgb,currentColor 16%,transparent);
--panel:color-mix(in srgb,currentColor 6%,transparent);}
.slide-wrap .frag>.fitwrap{flex:1 1 auto;min-height:0;width:100%;
display:flex;flex-direction:column;justify-content:inherit;align-items:inherit;
transform-origin:center center;}
.frag--dark{background:#0b0a14;color:#f4f3fb;background-image:
radial-gradient(120% 90% at 100% 0%,rgba(124,92,255,.22),transparent 55%),
radial-gradient(90% 80% at 100% 100%,rgba(227,0,140,.16),transparent 60%),
linear-gradient(180deg,#0b0a14,#0e0b1c);}
.frag--light{background:#f6f4fb;color:#161427;background-image:
radial-gradient(120% 90% at 100% 0%,rgba(124,92,255,.14),transparent 55%),
radial-gradient(90% 80% at 100% 100%,rgba(227,0,140,.08),transparent 60%);}
.frag--plain{background:var(--card);color:var(--fg);border:1px solid var(--border);}
.frag h1.ftitle{margin:0;font-weight:800;line-height:1.12;
font-size:clamp(27px,3.5vw,44px);letter-spacing:-.01em;
text-wrap:balance;overflow-wrap:break-word;word-break:keep-all;max-width:18ch;
margin-inline:auto;}
.frag:not(.frag--center) h1.ftitle{margin-inline:0;max-width:30ch;}
.frag .fsub{margin:14px 0 0;font-weight:700;font-size:clamp(16px,1.7vw,22px);
color:var(--msft-purple);line-height:1.3;}
/* Content slides: one consistent left-aligned header band (kicker stays top-right).
   Keeps every title starting at the same x so the deck reads evenly, never one-sided. */
.frag:not(.frag--center) .fhead{text-align:left;margin:0 0 6px;}
.frag:not(.frag--center) .fhead>.ftitle{margin-inline:0;max-width:30ch;}
.frag:not(.frag--center) .fhead>.flead,
.frag:not(.frag--center) .fhead>.fsub{margin-left:0;margin-right:0;max-width:74ch;}
.frag .feyebrow{font-weight:700;font-size:13px;letter-spacing:.08em;
text-transform:uppercase;color:var(--msft-purple);}
.frag .fkicker{position:absolute;top:18px;right:22px;font-size:12px;
color:var(--ink-dim);font-variant-numeric:tabular-nums;}
.frag p{line-height:1.62;}
.frag .flead{font-size:clamp(16px,1.6vw,20px);color:var(--ink-dim);line-height:1.62;}
.frag--center{justify-content:center;}
.frag--center.tc{align-items:center;text-align:center;}
.frag .fhead{margin:0 0 6px;}
.frag .fgrid{display:grid;gap:clamp(14px,1.8vw,26px);margin-top:auto;margin-bottom:auto;}
.frag .fgrid.c2{grid-template-columns:repeat(2,1fr);}
.frag .fgrid.imgL{grid-template-columns:1.55fr 1fr;}
.frag .fgrid.imgR{grid-template-columns:1fr 1.55fr;}
.frag .fgrid.c3{grid-template-columns:repeat(3,1fr);}
.frag .fgrid.c4{grid-template-columns:repeat(4,1fr);}
@media(max-width:760px){.frag .fgrid.c2,.frag .fgrid.imgL,.frag .fgrid.imgR,
.frag .fgrid.c3,.frag .fgrid.c4{grid-template-columns:1fr;}}
.frag .fcol{display:flex;flex-direction:column;align-items:center;text-align:center;gap:10px;}
.frag .fcol .h{font-weight:700;font-size:clamp(15px,1.5vw,18px);}
.frag .fcol .d{color:var(--ink-dim);font-size:clamp(13px,1.3vw,15.5px);line-height:1.55;}
.frag .down{width:34px;height:34px;border-radius:50%;flex:none;display:grid;place-items:center;
color:#fff;background:linear-gradient(135deg,var(--msft-blue),var(--msft-pink));
box-shadow:0 4px 14px rgba(0,0,0,.25);font-size:16px;}
.frag .fpanel{background:var(--panel);border:1px solid var(--hair);border-radius:14px;
padding:clamp(16px,1.8vw,22px);}
.frag .fpanel h3{margin:0 0 8px;font-size:clamp(16px,1.6vw,19px);}
.frag .fpanel p{margin:6px 0 0;color:var(--ink-dim);font-size:clamp(13px,1.3vw,15.5px);}
.frag .fstat-line{margin-top:auto;padding-top:18px;text-align:center;font-style:italic;
font-size:clamp(15px,1.6vw,20px);color:var(--ink-dim);line-height:1.5;}
.frag .fstat-line b{color:var(--msft-pink);font-style:normal;font-weight:800;}
.frag .fstats{display:flex;flex-wrap:wrap;gap:clamp(16px,2.4vw,44px);justify-content:center;
margin:auto 0;}
.frag .fstats .v{font-size:clamp(30px,4vw,52px);font-weight:800;line-height:1;
background:linear-gradient(135deg,var(--msft-blue),var(--msft-pink));
-webkit-background-clip:text;background-clip:text;color:transparent;}
.frag .fstats .l{margin-top:8px;color:var(--ink-dim);font-size:14px;text-align:center;}
.frag .fflow{display:flex;flex-wrap:wrap;align-items:center;gap:10px;}
.frag .fflow .node{background:var(--panel);border:1px solid var(--hair);border-radius:12px;
padding:12px 16px;font-weight:600;}
.frag .fflow .arr{color:var(--msft-purple);font-weight:800;font-size:1.2em;}
.frag ul.fbul{list-style:none;margin:14px 0 0;padding:0;}
.frag ul.fbul>li{position:relative;padding-left:1.5em;margin:11px 0;line-height:1.55;}
.frag ul.fbul>li::before{content:"";position:absolute;left:0;top:.55em;width:.55em;height:.55em;
border-radius:50%;background:linear-gradient(135deg,var(--msft-blue),var(--msft-pink));}
.frag .chips{display:flex;flex-wrap:wrap;gap:8px;}
.frag .chip{border:1px solid var(--hair);border-radius:999px;padding:6px 14px;font-size:14px;
background:var(--panel);}
.frag figure.video{margin:18px auto 0;width:100%;}
.frag .brandbar{display:flex;align-items:center;gap:8px;font-weight:600;font-size:14px;
color:var(--ink-dim);margin-bottom:18px;}
.frag .brandbar .sq{display:grid;grid-template-columns:1fr 1fr;width:18px;height:18px;gap:2px;}
.frag .brandbar .sq i{display:block;}
.frag .brandbar .sq i:nth-child(1){background:#f25022;}
.frag .brandbar .sq i:nth-child(2){background:#7fba00;}
.frag .brandbar .sq i:nth-child(3){background:#00a4ef;}
.frag .brandbar .sq i:nth-child(4){background:#ffb900;}
/* 실제 PPT 캡처 이미지 프레임 */
.frag figure.fshot{margin:0;width:100%;display:flex;flex-direction:column;justify-content:center;}
.frag figure.fshot img{display:block;width:100%;height:auto;border-radius:12px;
border:1px solid var(--hair);box-shadow:0 10px 34px rgba(0,0,0,.28);
background:#fff;}
.frag figure.fshot figcaption{margin-top:10px;color:var(--ink-dim);font-size:13px;text-align:center;}
.frag .fshot--zoom img{cursor:zoom-in;transition:transform .35s cubic-bezier(.2,.7,.2,1),box-shadow .35s;}
.frag .fshot--zoom img:hover{transform:scale(1.02);box-shadow:0 18px 50px rgba(0,0,0,.4);}
/* 이미지 그리드 안의 스크린샷은 셀을 가득 채워 충분히 크게 */
.frag .fgrid.imgL{grid-template-columns:1.7fr 1fr;align-items:stretch;}
.frag .fgrid.imgR{grid-template-columns:1fr 1.7fr;align-items:stretch;}
.frag .fgrid.imgL>figure.fshot,.frag .fgrid.imgR>figure.fshot{align-self:center;}
.frag .fgrid.imgL>figure.fshot img,.frag .fgrid.imgR>figure.fshot img{
width:100%;max-height:none;}
/* 이미지 비중이 작은 그리드는 더 넓게 보이도록 좌우 여백 축소 */
.frag:has(.fgrid.imgL) ,.frag:has(.fgrid.imgR){padding-left:clamp(20px,2.4vw,40px);padding-right:clamp(20px,2.4vw,40px);}
/* 단독(전폭) 스크린샷은 과하게 커지지 않도록 가운데 정렬 + 상한 */
.frag>figure.fshot,.frag .reveal>figure.fshot{max-width:1040px;margin-inline:auto;}
/* 이미지·다이어그램 라이트박스(클릭 시 확대) */
.lightbox{position:fixed;inset:0;z-index:50;background:rgba(0,0,0,.9);display:none;
align-items:center;justify-content:center;padding:4vh 4vw;cursor:zoom-out;}
.lightbox.on{display:flex;}
.lightbox img{max-width:100%;max-height:92vh;border-radius:10px;box-shadow:0 24px 70px rgba(0,0,0,.6);}
.lightbox .lbcontent{display:flex;align-items:center;justify-content:center;max-width:96vw;max-height:92vh;}
.lightbox .lbcontent svg{width:auto;height:auto;max-width:94vw;max-height:88vh;
background:#11131c;border-radius:12px;padding:22px;box-shadow:0 24px 70px rgba(0,0,0,.6);}
.lightbox .lbx{position:fixed;top:18px;right:22px;font-size:30px;color:#fff;opacity:.8;line-height:1;}
/* Mermaid 렌더 실패 폴백 */
.frag .mermaid.mmfail{display:block;}
.frag .mmfail-note{padding:16px 18px;border:1px dashed var(--hair);border-radius:12px;
color:var(--ink-dim);font-size:14px;text-align:center;background:color-mix(in srgb,currentColor 5%,transparent);}
/* 호버 인터랙션 */
.frag .fpanel{transition:transform .3s cubic-bezier(.2,.7,.2,1),border-color .3s,box-shadow .3s;}
.frag .fpanel:hover{transform:translateY(-4px);border-color:color-mix(in srgb,var(--msft-purple) 55%,transparent);
box-shadow:0 14px 36px rgba(0,0,0,.22);}
.frag .fcol{transition:transform .3s;}
.frag .fcol:hover{transform:translateY(-3px);}
.frag .chip{transition:transform .2s,background .2s,border-color .2s;}
.frag .chip:hover{transform:translateY(-2px);border-color:color-mix(in srgb,var(--msft-blue) 60%,transparent);}
.frag .fflow .node{transition:transform .25s,box-shadow .25s;}
.frag .fflow .node:hover{transform:translateY(-2px);box-shadow:0 8px 22px rgba(0,0,0,.22);}
/* 스크롤 진입 애니메이션 */
.slide-wrap{opacity:0;transform:translateY(28px);
transition:opacity .7s ease,transform .7s cubic-bezier(.2,.7,.2,1);will-change:opacity,transform;}
.slide-wrap.in{opacity:1;transform:none;}
.slide-wrap.in .reveal{animation:fadeUp .6s both;}
.slide-wrap.in .reveal:nth-child(2){animation-delay:.08s;}
.slide-wrap.in .reveal:nth-child(3){animation-delay:.16s;}
.slide-wrap.in .reveal:nth-child(4){animation-delay:.24s;}
.slide-wrap.in .reveal:nth-child(5){animation-delay:.32s;}
.slide-wrap.in .reveal:nth-child(6){animation-delay:.40s;}
@keyframes fadeUp{from{opacity:0;transform:translateY(16px);}to{opacity:1;transform:none;}}
/* 다이어그램 라이브러리(Mermaid) */
.frag .fdiagram{margin:18px auto 0;max-width:1100px;width:100%;position:relative;}
.frag .fdiagram .mermaid{display:flex;justify-content:center;line-height:1.3;}
.frag .fdiagram .mermaid svg{max-width:100%;height:auto;}
.frag .fdiagram.zoomable .mermaid{cursor:zoom-in;}
.frag .fdiagram .fdiagram-zoom{position:absolute;top:6px;right:6px;z-index:2;display:inline-flex;align-items:center;gap:5px;
padding:5px 11px;font-size:12px;font-weight:600;border-radius:999px;cursor:zoom-in;
background:color-mix(in srgb,var(--msft-purple) 22%,#11131c);
border:1px solid color-mix(in srgb,var(--msft-purple) 50%,transparent);color:#e8eaf2;opacity:.8;
transition:opacity .2s,transform .2s,box-shadow .2s;}
.frag .fdiagram .fdiagram-zoom:hover{opacity:1;transform:translateY(-1px);box-shadow:0 6px 18px rgba(0,0,0,.35);}
.frag .fdiagram .fdiagram-cap{margin-top:10px;text-align:center;color:var(--ink-dim);font-size:13px;}
/* 방사형(orbit) 다이어그램 */
.frag .forbit{position:relative;width:min(560px,82vw);aspect-ratio:1;margin:22px auto 0;}
.frag .forbit .ring{position:absolute;inset:14%;border:1.5px dashed color-mix(in srgb,currentColor 22%,transparent);
border-radius:50%;animation:spin 60s linear infinite;}
.frag .forbit .ring::before{content:"";position:absolute;inset:11%;border:1px dashed color-mix(in srgb,currentColor 12%,transparent);border-radius:50%;}
.frag .forbit .ocenter{position:absolute;left:50%;top:50%;transform:translate(-50%,-50%);
width:38%;aspect-ratio:1;border-radius:50%;display:grid;place-items:center;text-align:center;gap:6px;
background:radial-gradient(circle at 35% 30%,rgba(139,92,246,.28),rgba(20,22,34,.85));
border:1px solid color-mix(in srgb,currentColor 18%,transparent);box-shadow:0 18px 50px rgba(0,0,0,.45);padding:8px;}
.frag .forbit .ocenter img{width:46%;height:auto;filter:drop-shadow(0 6px 18px rgba(109,92,246,.5));}
.frag .forbit .ocenter b{font-size:clamp(12px,1.3vw,15px);line-height:1.2;}
.frag .forbit .onode{position:absolute;transform:translate(-50%,-50%);width:31%;text-align:center;
display:flex;flex-direction:column;align-items:center;gap:6px;}
.frag .forbit .onode .ic{width:clamp(40px,5vw,56px);aspect-ratio:1;border-radius:50%;display:grid;place-items:center;
font-size:clamp(18px,2.4vw,24px);box-shadow:0 8px 22px rgba(0,0,0,.35);transition:transform .25s,box-shadow .25s;}
.frag .forbit .onode:hover .ic{transform:translateY(-3px) scale(1.06);box-shadow:0 12px 30px rgba(0,0,0,.45);}
.frag .forbit .onode .ot{font-weight:700;font-size:clamp(12px,1.35vw,15px);line-height:1.2;}
.frag .forbit .onode .od{font-size:clamp(10px,1.1vw,12px);color:var(--ink-dim);line-height:1.3;}
@keyframes spin{to{transform:rotate(360deg);}}
/* === 하단 컨트롤 바 === */
.deckbar{position:fixed;left:50%;bottom:18px;transform:translateX(-50%);z-index:45;
display:flex;align-items:center;gap:4px;padding:6px 8px;border-radius:999px;
background:color-mix(in srgb,var(--card) 86%,transparent);backdrop-filter:blur(12px);
-webkit-backdrop-filter:blur(12px);border:1px solid var(--border);
box-shadow:0 10px 34px rgba(0,0,0,.3);font-size:13px;color:var(--fg);
opacity:.32;transition:opacity .25s,transform .25s;}
.deckbar:hover,.deckbar:focus-within{opacity:1;}
.deckbar button{display:inline-flex;align-items:center;justify-content:center;gap:6px;
min-width:34px;height:34px;padding:0 12px;border-radius:999px;border:1px solid transparent;
background:none;color:inherit;cursor:pointer;font-size:15px;line-height:1;}
.deckbar button:hover{background:var(--active);}
.deckbar .dbcount{min-width:62px;text-align:center;font-variant-numeric:tabular-nums;
font-weight:600;cursor:default;}
.deckbar .dbcount:hover{background:none;}
.deckbar .sep{width:1px;height:20px;background:var(--border);margin:0 3px;}
.deckbar .lbl{font-size:13px;font-weight:600;}
@media (max-width:560px){.deckbar .lbl{display:none;}.deckbar button{padding:0 8px;}}
@media print{.deckbar{display:none;}}
/* === 발표(전체화면 단일 슬라이드) 오버레이 === */
.present{position:fixed;inset:0;z-index:60;background:#05060a;display:none;
align-items:center;justify-content:center;}
body.presenting{overflow:hidden;}
body.presenting .present{display:flex;}
.present .stagebox{position:relative;width:min(96vw,calc(94vh*16/9));
height:min(94vh,calc(96vw*9/16));overflow:hidden;border-radius:14px;
background:var(--bg);box-shadow:0 30px 90px rgba(0,0,0,.6);}
.fit-holder{position:absolute;top:0;left:0;transform-origin:top left;}
.fit-holder .slide-wrap{margin:0;max-width:none;opacity:1;transform:none;}
.fit-holder .frag{min-height:0;}
.present .pbtn,.present .pexit{position:absolute;z-index:3;border-radius:50%;
border:1px solid rgba(255,255,255,.2);background:rgba(255,255,255,.08);color:#fff;
cursor:pointer;display:grid;place-items:center;opacity:.45;transition:opacity .2s,background .2s;}
.present .pbtn:hover,.present .pexit:hover{opacity:1;background:rgba(255,255,255,.18);}
.present .pbtn{top:50%;transform:translateY(-50%);width:48px;height:48px;font-size:24px;}
.present .pprev{left:20px;} .present .pnext{right:20px;}
.present .pexit{top:18px;right:20px;width:40px;height:40px;font-size:18px;}
.present .pnum{position:absolute;bottom:16px;left:50%;transform:translateX(-50%);
color:#fff;opacity:.55;font-size:14px;font-variant-numeric:tabular-nums;z-index:3;}
/* === 발표자 모드(별도 창) === */
.presenter-room{display:none;}
body.presenter{overflow:hidden;}
body.presenter .layout,body.presenter .deckbar,body.presenter .topbar,
body.presenter .side-open,body.presenter .present{display:none!important;}
body.presenter .presenter-room{display:grid;position:fixed;inset:0;background:#0b0d14;
color:#e8eaf2;gap:14px;padding:16px;
grid-template-columns:1.55fr 1fr;grid-template-rows:1.1fr 1.4fr auto;
grid-template-areas:"cur next" "cur script" "bar bar";}
.presenter-room .plabel{font-size:12px;font-weight:700;letter-spacing:.08em;
text-transform:uppercase;color:#8b93b4;margin-bottom:8px;}
.presenter-room .pcur{grid-area:cur;display:flex;flex-direction:column;min-height:0;}
.presenter-room .pnext{grid-area:next;display:flex;flex-direction:column;min-height:0;}
.presenter-room .pscript{grid-area:script;display:flex;flex-direction:column;min-height:0;
background:#11141e;border:1px solid #232838;border-radius:12px;padding:16px 18px;}
.presenter-room .pscript-body{overflow:auto;white-space:pre-wrap;line-height:1.62;
font-size:clamp(16px,1.5vw,22px);color:#dfe3f0;}
.presenter-room .pscript-body.empty{color:#6b7390;font-style:italic;}
.presenter-room .pcur .stagebox,.presenter-room .pnext .stagebox{position:relative;
flex:1;min-height:0;overflow:hidden;border-radius:12px;background:#05060a;
border:1px solid #232838;}
.presenter-room .pbar{grid-area:bar;display:flex;align-items:center;gap:14px;
background:#11141e;border:1px solid #232838;border-radius:12px;padding:10px 16px;}
.presenter-room .pbar .ptime{font-size:22px;font-weight:700;font-variant-numeric:tabular-nums;
letter-spacing:.02em;}
.presenter-room .pbar .pcounter{font-size:15px;color:#aeb6d4;font-variant-numeric:tabular-nums;}
.presenter-room .pbar .grow{flex:1;}
.presenter-room .pbar button{height:40px;min-width:44px;padding:0 16px;border-radius:10px;
border:1px solid #2c3247;background:#1a1f2e;color:#e8eaf2;cursor:pointer;font-size:15px;
font-weight:600;display:inline-flex;align-items:center;gap:6px;}
.presenter-room .pbar button:hover{background:#252b3d;}
.presenter-room .pbar button.accent{background:linear-gradient(135deg,#3b82f6,#8b5cf6);border:none;}
@media (max-width:780px){body.presenter .presenter-room{grid-template-columns:1fr;
grid-template-rows:auto auto 1fr auto;grid-template-areas:"cur" "next" "script" "bar";}}
@media (prefers-reduced-motion:reduce){
.slide-wrap{opacity:1;transform:none;transition:none;}
.slide-wrap.in .reveal{animation:none;}
.frag .fpanel:hover,.frag .fcol:hover,.frag .chip:hover,.frag .fflow .node:hover,
.frag .fshot--zoom img:hover{transform:none;}
.frag .forbit .ring{animation:none;}
}
"""


# 발표/발표자/전체화면 인터랙션 (f-string 아님 — 중괄호 그대로 사용 가능)
PRESENT_JS = r"""
(function(){
  var NOTES = window.__NOTES || {};
  var DECK  = window.__DECK  || 'deck';
  var slides = [].slice.call(document.querySelectorAll('.slide-wrap'));
  var total = slides.length;
  if(!total) return;

  // === 모든 슬라이드 16:9 고정 + 내용 자동 맞춤(넘치면 축소) ===
  function buildFit(){
    document.querySelectorAll('.slide-wrap > .frag').forEach(function(frag){
      if(frag.querySelector(':scope > .fitwrap')) return;
      var w=document.createElement('div'); w.className='fitwrap';
      while(frag.firstChild) w.appendChild(frag.firstChild);
      frag.appendChild(w);
    });
  }
  function fitFrag(frag){
    var w=frag.querySelector(':scope > .fitwrap'); if(!w) return;
    w.style.transform='none';
    var cs=getComputedStyle(frag);
    var ah=frag.clientHeight-parseFloat(cs.paddingTop)-parseFloat(cs.paddingBottom);
    var aw=frag.clientWidth-parseFloat(cs.paddingLeft)-parseFloat(cs.paddingRight);
    if(ah<=0||aw<=0) return;
    var s=Math.min(1, ah/w.scrollHeight, aw/w.scrollWidth);
    if(isFinite(s)&&s>0&&s<0.999) w.style.transform='scale('+s+')';
  }
  function fitAll(){ document.querySelectorAll('.slide-wrap > .frag').forEach(fitFrag); }
  window.__fitAll=fitAll;
  buildFit(); fitAll();
  var _ft; function fitSoon(){ clearTimeout(_ft); _ft=setTimeout(fitAll,60); }
  window.addEventListener('load', fitAll);
  window.addEventListener('resize', fitSoon);
  if(document.fonts&&document.fonts.ready) document.fonts.ready.then(fitAll);
  document.querySelectorAll('.slide-wrap img').forEach(function(im){
    if(!im.complete) im.addEventListener('load', fitSoon, {once:true});
  });

  var state = { index: 0 };
  var isPresenter = /present=presenter/.test(location.hash);

  var chan = null;
  try { chan = new BroadcastChannel('deck:'+DECK); } catch(e) { chan = null; }
  function send(msg){ if(chan){ try{ chan.postMessage(msg); }catch(e){} } }

  function currentScrollIndex(){
    var mid=window.innerHeight*0.5, best=0, bestD=Infinity;
    for(var i=0;i<slides.length;i++){
      var r=slides[i].getBoundingClientRect();
      if(r.top<=mid && r.bottom>=mid) return i;
      var c=(r.top+r.bottom)/2, d=Math.abs(c-mid);
      if(d<bestD){ bestD=d; best=i; }
    }
    return best;
  }

  var DESIGN_W = 1280;
  function fitInto(box, slideEl){
    if(!box) return;
    box.innerHTML='';
    if(!slideEl) return;
    var holder=document.createElement('div');
    holder.className='fit-holder';
    holder.style.width=DESIGN_W+'px';
    var node=slideEl.cloneNode(true);
    node.classList.add('in');
    holder.appendChild(node);
    box.appendChild(holder);
    requestAnimationFrame(function(){
      var bw=box.clientWidth, bh=box.clientHeight;
      var nh=holder.scrollHeight||720, nw=DESIGN_W;
      var s=Math.min(bw/nw, bh/nh);
      holder.style.transform='scale('+s+')';
      holder.style.left=Math.max(0,(bw-nw*s)/2)+'px';
      holder.style.top=Math.max(0,(bh-nh*s)/2)+'px';
      var cf=holder.querySelector('.slide-wrap > .frag'); if(cf) fitFrag(cf);
    });
  }

  var present=document.querySelector('.present');
  var pStage=present? present.querySelector('.stagebox'):null;
  var pNum=present? present.querySelector('.pnum'):null;
  var presenterLinked=false;
  function requestFS(){
    try{ if(!document.fullscreenElement && document.documentElement.requestFullscreen)
      document.documentElement.requestFullscreen(); }catch(e){}
  }
  function renderPresent(){
    fitInto(pStage, slides[state.index]);
    if(pNum) pNum.textContent=(state.index+1)+' / '+total;
  }
  function enterPresent(){
    presenterLinked=false;
    document.body.classList.add('presenting');
    renderPresent();
    requestFS();
  }
  function exitPresent(){
    presenterLinked=false;
    document.body.classList.remove('presenting');
    if(pStage) pStage.innerHTML='';
    try{ if(document.fullscreenElement && document.exitFullscreen) document.exitFullscreen(); }catch(e){}
    if(slides[state.index]) slides[state.index].scrollIntoView({block:'start'});
  }

  var room=document.querySelector('.presenter-room');
  var curBox=room? room.querySelector('.pcur-box'):null;
  var nextBox=room? room.querySelector('.pnext-box'):null;
  var scriptBody=room? room.querySelector('.pscript-body'):null;
  var counter=room? room.querySelector('.pcounter'):null;
  var timeEl=room? room.querySelector('.ptime'):null;
  function renderPresenter(){
    fitInto(curBox, slides[state.index]);
    fitInto(nextBox, slides[state.index+1]||null);
    var n=NOTES[String(state.index+1)]||'';
    if(scriptBody){
      scriptBody.textContent=n||'이 슬라이드에는 발표자 노트가 없습니다.';
      scriptBody.classList.toggle('empty', !n);
    }
    if(counter) counter.textContent=(state.index+1)+' / '+total;
  }

  var bar=document.querySelector('.deckbar');
  var barCount=bar? bar.querySelector('.dbcount'):null;
  function updateBar(){ if(barCount) barCount.textContent=(state.index+1)+' / '+total; }

  var navAt=0;
  function go(i, opts){
    opts=opts||{};
    i=Math.max(0,Math.min(total-1,i));
    state.index=i;
    updateBar();
    if(document.body.classList.contains('presenter')) renderPresenter();
    else if(document.body.classList.contains('presenting')) renderPresent();
    else if(opts.scroll!==false && slides[i]){
      navAt=Date.now();
      slides[i].scrollIntoView({behavior:opts.smooth===false?'auto':'smooth',block:'start'});
    }
    if(!opts.silent) send({t:'goto', i:state.index});
  }
  function next(){ go(state.index+1); }
  function prev(){ go(state.index-1); }

  function toggleFS(){
    var d=document;
    if(!d.fullscreenElement){ if(d.documentElement.requestFullscreen) d.documentElement.requestFullscreen(); }
    else if(d.exitFullscreen){ d.exitFullscreen(); }
  }

  function openPresenter(){
    presenterLinked=true;
    document.body.classList.add('presenting');
    renderPresent();
    var url=location.pathname+location.search+'#present=presenter';
    var w=window.open(url,'deckPresenter_'+DECK,'width=1280,height=820');
    if(!w){ alert('팝업이 차단되었습니다. 발표자 창을 열려면 팝업을 허용해 주세요.'); presenterLinked=false; return; }
    /* The popup steals OS focus the instant it opens, and Chrome drops the
       opener's fullscreen as soon as it loses focus — that is why the audience
       window flashes fullscreen then snaps back. Push focus off the popup and
       back onto the audience window, THEN request fullscreen, all inside this
       same click gesture so the request still counts as user-initiated. */
    try{ w.blur(); }catch(e){}
    try{ window.focus(); }catch(e){}
    requestFS();
    /* Belt-and-suspenders: if the popup grabs focus again a tick later, refocus
       the audience window so the OS keeps it fullscreen. (Cannot re-request
       fullscreen here — no gesture — but holding focus prevents the auto-exit.) */
    setTimeout(function(){ try{ w.blur(); }catch(e){} try{ window.focus(); }catch(e){} }, 250);
  }

  if(chan){
    chan.onmessage=function(ev){
      var m=ev.data||{};
      if(m.t==='goto'){ go(m.i,{silent:true,scroll:true}); }
      else if(m.t==='hello'){ send({t:'goto', i:state.index}); }
    };
  }

  document.addEventListener('keydown',function(e){
    if(e.altKey||e.ctrlKey||e.metaKey) return;
    var tag=(e.target&&e.target.tagName)||'';
    if(tag==='SELECT'||tag==='INPUT'||tag==='TEXTAREA') return;
    var mode = document.body.classList.contains('presenter') ? 'presenter'
             : document.body.classList.contains('presenting') ? 'present' : 'scroll';
    if(e.key==='ArrowRight'||e.key==='PageDown'){ next(); if(mode!=='scroll') e.preventDefault(); }
    else if(e.key==='ArrowLeft'||e.key==='PageUp'){ prev(); if(mode!=='scroll') e.preventDefault(); }
    else if(e.key===' '){ if(mode!=='scroll'){ next(); e.preventDefault(); } }
    else if(e.key==='Escape'){ if(mode==='present') exitPresent(); }
    else if(e.key==='f'||e.key==='F'){ if(mode!=='scroll') toggleFS(); }
  });

  if(isPresenter){
    document.body.classList.add('presenter');
    var t0=Date.now();
    function tick(){ if(!timeEl) return; var s=Math.floor((Date.now()-t0)/1000);
      var mm=String(Math.floor(s/60)).padStart(2,'0'), ss=String(s%60).padStart(2,'0');
      timeEl.textContent=mm+':'+ss; }
    setInterval(tick,1000); tick();
    if(room){
      var rp=room.querySelector('.pp-prev'), rn=room.querySelector('.pp-next'),
          rf=room.querySelector('.pp-fs'), rr=room.querySelector('.pp-reset');
      if(rp) rp.onclick=prev; if(rn) rn.onclick=next; if(rf) rf.onclick=toggleFS;
      if(rr) rr.onclick=function(){ t0=Date.now(); tick(); };
    }
    send({t:'hello'});
    renderPresenter();
    window.addEventListener('load', renderPresenter);
    setTimeout(renderPresenter, 1500);
    window.addEventListener('resize', renderPresenter);
    return;
  }

  if(bar){
    var b;
    if((b=bar.querySelector('.db-prev'))) b.onclick=prev;
    if((b=bar.querySelector('.db-next'))) b.onclick=next;
    if((b=bar.querySelector('.db-fs'))) b.onclick=toggleFS;
    if((b=bar.querySelector('.db-present'))) b.onclick=enterPresent;
    if((b=bar.querySelector('.db-presenter'))) b.onclick=openPresenter;
  }
  if(present){
    present.querySelector('.pprev').onclick=prev;
    present.querySelector('.pnext').onclick=next;
    present.querySelector('.pexit').onclick=exitPresent;
  }
  var ticking=false;
  window.addEventListener('scroll',function(){
    if(document.body.classList.contains('presenting')) return;
    if(Date.now()-navAt < 700) return;
    if(ticking) return; ticking=true;
    requestAnimationFrame(function(){ state.index=currentScrollIndex(); updateBar(); ticking=false; });
  },{passive:true});
  window.addEventListener('resize',function(){
    if(document.body.classList.contains('presenting')) renderPresent();
  });
  document.addEventListener('fullscreenchange',function(){
    if(document.fullscreenElement) return;
    /* Plain Present mode: exiting fullscreen exits the view. But when a
       presenter popup is linked, the popup opening transiently drops the
       opener's fullscreen — keep the single-slide audience view (user exits
       via Esc / the exit button instead of being kicked back to scroll). */
    if(!presenterLinked && document.body.classList.contains('presenting')) exitPresent();
  });
  state.index=currentScrollIndex();
  updateBar();
})();
"""


def _block_html(b: dict, title_norm: str = "") -> str:
    if b["type"] == "bullets":
        lis = []
        for it in b["items"]:
            if title_norm and _norm(it["text"]) == title_norm:
                continue
            cls = f"lv{min(it['level'], 4)}"
            if it["bold"]:
                cls += " b"
            align = f' style="text-align:{it["align"]}"' if it["align"] != "left" else ""
            lis.append(f'<li class="{cls}"{align}>{escape(it["text"])}</li>')
        return f'<ul class="bullets">{"".join(lis)}</ul>' if lis else ""
    if b["type"] == "table":
        rows = "".join(
            "<tr>" + "".join(f"<td>{escape(c)}</td>" for c in row) + "</tr>"
            for row in b["rows"])
        return f'<div class="tbl"><table>{rows}</table></div>'
    if b["type"] == "image":
        return f'<figure class="img"><img src="{b["src"]}" alt="" loading="lazy"></figure>'
    if b["type"] in ("video", "embed"):
        return _video_html(b)
    return ""


def render_present(slide: dict, title: str) -> str:
    """슬라이드 내용을 발표용 카드 본문(불릿/표/이미지)으로 렌더."""
    title_norm = _norm(title)
    parts = [h for b in slide["blocks"] if (h := _block_html(b, title_norm))]
    if not parts:
        parts.append('<p class="empty">(내용 없음)</p>')
    return f'<div class="card-body">{"".join(parts)}</div>'


def slugify(name: str) -> str:
    s = re.sub(r"\s+", "-", name.strip())
    s = re.sub(r"[^0-9A-Za-z가-힣\-_]", "", s)
    s = re.sub(r"-{2,}", "-", s).strip("-")
    return s or "deck"


# --------------------------------------------------------------------------- #
# 에이전트 재구성 콘텐츠 (content.json) — 슬라이드 내용을 고객 전달용으로 재작성
# --------------------------------------------------------------------------- #
def load_content(*candidates: Path | None) -> dict:
    """첫 번째로 존재하는 content.json 을 {슬라이드번호(str): entry} 로."""
    for c in candidates:
        if c and c.exists():
            try:
                data = json.loads(c.read_text(encoding="utf-8"))
                slides = data.get("slides", data) if isinstance(data, dict) else {}
                return {str(k): v for k, v in slides.items()}
            except Exception:
                return {}
    return {}


def _etext(s: str) -> str:
    return escape(str(s)).replace("\n", "<br>")


def _blocks_html(blocks: list) -> str:
    out = []
    for b in blocks or []:
        t = b.get("type")
        if t in ("lead", "para"):
            cls = ' class="lead"' if t == "lead" else ""
            out.append(f"<p{cls}>{_etext(b['text'])}</p>")
        elif t == "bullets":
            lis = ""
            for it in b["items"]:
                if isinstance(it, dict):
                    sub = "".join(f"<li>{_etext(s)}</li>" for s in it.get("sub", []))
                    sub = f"<ul>{sub}</ul>" if sub else ""
                    lis += f"<li>{_etext(it['text'])}{sub}</li>"
                else:
                    lis += f"<li>{_etext(it)}</li>"
            out.append(f'<ul class="bul">{lis}</ul>')
        elif t == "cards":
            cs = ""
            for c in b["items"]:
                d = f"<p>{_etext(c['desc'])}</p>" if c.get("desc") else ""
                cs += f"<div class=\"gc\"><h4>{_etext(c['title'])}</h4>{d}</div>"
            out.append(f'<div class="grid">{cs}</div>')
        elif t == "steps":
            ss = ""
            for i, c in enumerate(b["items"], 1):
                d = f"<p>{_etext(c['desc'])}</p>" if c.get("desc") else ""
                ss += (f'<li><span class="sn">{i}</span>'
                       f"<div><h4>{_etext(c['title'])}</h4>{d}</div></li>")
            out.append(f'<ol class="steps">{ss}</ol>')
        elif t == "flow":
            arrow = ' <span class="fa">&rarr;</span> '
            out.append('<div class="flow">' +
                       arrow.join(f'<span class="fi">{_etext(x)}</span>' for x in b["items"]) +
                       "</div>")
        elif t == "stats":
            ss = "".join(
                f'<div class="stat"><div class="sv">{_etext(c["value"])}</div>'
                f'<div class="sl">{_etext(c["label"])}</div></div>' for c in b["items"])
            out.append(f'<div class="stats">{ss}</div>')
        elif t == "callout":
            style = b.get("style", "info")
            head = f"<strong>{_etext(b['title'])}</strong> " if b.get("title") else ""
            out.append(f'<div class="callout {style}">{head}{_etext(b["text"])}</div>')
        elif t == "table":
            head = ""
            if b.get("headers"):
                head = "<tr>" + "".join(f"<th>{_etext(h)}</th>" for h in b["headers"]) + "</tr>"
            rows = "".join(
                "<tr>" + "".join(f"<td>{_etext(c)}</td>" for c in r) + "</tr>"
                for r in b["rows"])
            out.append(f'<div class="tbl"><table>{head}{rows}</table></div>')
        elif t == "quote":
            cite = f"<cite>— {_etext(b['cite'])}</cite>" if b.get("cite") else ""
            out.append(f"<blockquote>{_etext(b['text'])}{cite}</blockquote>")
        elif t == "columns":
            cs = ""
            for col in b["items"]:
                cs += (f'<div class="col"><h4>{_etext(col["title"])}</h4>'
                       f'{_blocks_html(col.get("blocks", []))}</div>')
            out.append(f'<div class="cols">{cs}</div>')
        elif t == "image":
            src = b.get("src") or b.get("url")
            if src:
                cap = f"<figcaption>{_etext(b['caption'])}</figcaption>" if b.get("caption") else ""
                out.append(f'<figure class="img"><img src="{escape(src)}" alt="" loading="lazy">{cap}</figure>')
        elif t in ("video", "embed"):
            out.append(_video_html(b))
    return "".join(out)


def render_authored(entry: dict) -> tuple[str, str]:
    """content.json 의 한 슬라이드 entry → (본문 HTML, 카드 추가 클래스)."""
    parts = []
    if entry.get("lead"):
        parts.append(f'<p class="lead">{_etext(entry["lead"])}</p>')
    parts.append(_blocks_html(entry.get("blocks", [])))
    body = f'<div class="card-body">{"".join(parts)}</div>'
    extra = " card--hero" if entry.get("layout") in ("section", "title") else ""
    return body, extra


def load_fragments(*dirs: Path | None) -> dict:
    """모델이 작성한 슬라이드별 HTML 조각(NN.html) → {슬라이드번호(str): html}.
    여러 후보 폴더 중 처음 존재하는 폴더를 사용한다."""
    for d in dirs:
        if d and d.is_dir():
            frags = {}
            for f in d.glob("*.html"):
                m = re.match(r"0*(\d+)$", f.stem)
                if m:
                    frags[str(int(m.group(1)))] = f.read_text(encoding="utf-8")
            if frags:
                return frags
    return {}


def _slide_section_html(s: dict, title: str, total: int, content: dict,
                        fragments: dict | None = None) -> str:
    """한 슬라이드의 <section> HTML.
    우선순위: ① 모델 작성 HTML 조각 → ② content.json 재구성 → ③ 원문 불릿."""
    sid = f"slide-{s['num']}"
    fragments = fragments or {}
    frag = fragments.get(str(s["num"]))
    if frag is not None:
        return f'<div class="slide-wrap" id="{sid}">{frag}</div>'
    entry = content.get(str(s["num"]))
    if entry is not None:
        body, extra = render_authored(entry)
    else:
        body, extra = render_present(s, title), ""
    return (
        f'<div class="slide-wrap" id="{sid}"><article class="card{extra}">'
        f'<div class="card-head"><span class="kicker">{s["num"]} / {total}</span>'
        f'<h2>{escape(title)}</h2></div>'
        f"{body}</article></div>"
    )


def render_version_html(rel: dict, deck: list[dict], titles: dict, versions: list[dict],
                        content: dict | None = None, fragments: dict | None = None,
                        theme: dict | None = None) -> str:
    """한 버전(릴리즈)의 자체완결 HTML. PPT 폴더 안에 <release-id>.html 로 저장."""
    content = content or {}
    fragments = fragments or {}
    if theme is None:
        theme = extract_theme(rel["pptx"])
    theme_css = render_theme_css(theme)
    mermaid_vars = mermaid_theme_json(theme)
    deck_name = rel["stem"]
    notes_json = json.dumps({str(s["num"]): s.get("notes", "") for s in deck},
                            ensure_ascii=False)
    deck_json = json.dumps(deck_name, ensure_ascii=False)
    nav_items, sections = [], []
    for s in deck:
        title = resolve_title(s, titles)
        sid = f"slide-{s['num']}"
        nav_items.append(
            f'<li><a href="#{sid}" data-target="{sid}">'
            f'<span class="n">{s["num"]}</span><span>{escape(title)}</span></a></li>'
        )
        sections.append(_slide_section_html(s, title, len(deck), content, fragments))

    # 버전 전환기 (같은 PPT 폴더 내 다른 릴리즈로 이동)
    opts = []
    for v in versions:
        sel = " selected" if v["id"] == rel["id"] else ""
        opts.append(f'<option value="{v["id"]}.html"{sel}>{escape(v["ts_human"])}</option>')
    ver_switch = (
        '<div class="ver">버전<select onchange="location.href=this.value">'
        + "".join(opts)
        + "</select></div>"
    ) if len(versions) > 1 else f'<div class="ver">{escape(rel["ts_human"])}</div>'

    return f"""<!DOCTYPE html>
<html lang="ko"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(deck_name)}</title>
<style>{CSS}</style>
<style>{theme_css}</style></head><body>
<div class="topbar"><button onclick="document.querySelector('aside').classList.toggle('open')">&#9776;</button><span class="t">{escape(deck_name)}</span></div>
<button class="side-open" onclick="toggleSide()" title="사이드바 열기" aria-label="사이드바 열기">&#9776;</button>
<div class="layout">
<aside>
<button class="collapse-btn" onclick="toggleSide()" title="사이드바 접기" aria-label="사이드바 접기">&#10094;</button>
<div class="deck-title">{escape(deck_name)}</div>
{ver_switch}
<div class="ver"><a href="index.html" style="color:var(--accent)">버전 히스토리 {len(versions)}개 &rarr;</a></div>
<nav><ol>{''.join(nav_items)}</ol></nav>
<div class="ver" style="padding-top:16px"><a href="../index.html" style="color:var(--accent)">&larr; 전체 덱 목록</a></div>
</aside>
<main>{''.join(sections)}</main>
</div>
<div class="deckbar" role="group" aria-label="발표 제어">
<button class="db-prev" title="이전 슬라이드 (←)" aria-label="이전 슬라이드">&#8249;</button>
<span class="dbcount">1 / {len(deck)}</span>
<button class="db-next" title="다음 슬라이드 (→)" aria-label="다음 슬라이드">&#8250;</button>
<span class="sep"></span>
<button class="db-present" title="발표 시작 (전체화면 슬라이드쇼)"><span>&#9654;</span><span class="lbl">발표</span></button>
<button class="db-presenter" title="발표자 모드 (2화면: 슬라이드 + 발표자 노트)"><span>&#127908;</span><span class="lbl">발표자</span></button>
<button class="db-fs" title="전체화면 전환" aria-label="전체화면">&#9974;</button>
</div>
<div class="present" aria-hidden="true">
<button class="pbtn pprev" aria-label="이전">&#8249;</button>
<div class="stagebox"></div>
<button class="pbtn pnext" aria-label="다음">&#8250;</button>
<button class="pexit" aria-label="발표 종료 (Esc)">&#10005;</button>
<div class="pnum"></div>
</div>
<div class="presenter-room">
<div class="pcur"><div class="plabel">현재 슬라이드</div><div class="stagebox pcur-box"></div></div>
<div class="pnext"><div class="plabel">다음 슬라이드</div><div class="stagebox pnext-box"></div></div>
<div class="pscript"><div class="plabel">발표자 노트</div><div class="pscript-body"></div></div>
<div class="pbar">
<div class="ptime">00:00</div>
<button class="pp-reset" title="타이머 초기화">&#8635;</button>
<div class="pcounter">1 / {len(deck)}</div>
<div class="grow"></div>
<button class="pp-prev">&#8249; 이전</button>
<button class="pp-next accent">다음 &#8250;</button>
<button class="pp-fs" title="전체화면">&#9974;</button>
</div>
</div>
<script>
function toggleSide(){{document.body.classList.toggle('side-collapsed');}}
const links=[...document.querySelectorAll('nav a')];
const map=new Map(links.map(a=>[a.dataset.target,a]));
const obs=new IntersectionObserver((es)=>{{
 es.forEach(e=>{{if(e.isIntersecting){{
  links.forEach(a=>a.classList.remove('active'));
  const a=map.get(e.target.id); if(a){{a.classList.add('active');
   a.scrollIntoView({{block:'nearest'}});}}
 }}}});
}},{{rootMargin:'-45% 0px -50% 0px'}});
document.querySelectorAll('.slide-wrap').forEach(s=>obs.observe(s));
const rm=window.matchMedia('(prefers-reduced-motion:reduce)').matches;
if(rm){{document.querySelectorAll('.slide-wrap').forEach(s=>s.classList.add('in'));}}
else{{
 const rev=new IntersectionObserver((es)=>{{
  es.forEach(e=>{{if(e.isIntersecting){{e.target.classList.add('in');rev.unobserve(e.target);}}}});
 }},{{rootMargin:'0px 0px -12% 0px',threshold:.12}});
 document.querySelectorAll('.slide-wrap').forEach(s=>rev.observe(s));
}}
links.forEach(a=>a.addEventListener('click',()=>document.querySelector('aside').classList.remove('open')));
(function(){{
 const lb=document.createElement('div');lb.className='lightbox';
 lb.innerHTML='<span class="lbx">&times;</span><div class="lbcontent"></div>';
 document.body.appendChild(lb);
 const lbc=lb.querySelector('.lbcontent');
 const close=()=>{{lb.classList.remove('on');lbc.innerHTML='';}};
 window.__lb={{open(node){{lbc.innerHTML='';lbc.appendChild(node);lb.classList.add('on');}},close}};
 [...document.querySelectorAll('.fshot--zoom img')].forEach(im=>im.addEventListener('click',()=>{{
  const c=new Image();c.src=im.currentSrc||im.src;c.alt=im.alt||'';window.__lb.open(c);
 }}));
 lb.addEventListener('click',close);
 document.addEventListener('keydown',e=>{{if(e.key==='Escape')close();}});
}})();
</script>
<script type="module">
(async()=>{{
 const nodes=[...document.querySelectorAll('.mermaid')];
 if(!nodes.length) return;
 try{{
  const {{default:mermaid}}=await import('https://cdn.jsdelivr.net/npm/mermaid@11/dist/mermaid.esm.min.mjs');
  mermaid.initialize({{startOnLoad:false,suppressErrorRendering:true,securityLevel:'loose',theme:'base',
   fontFamily:'inherit',
   themeVariables:{mermaid_vars},
   flowchart:{{curve:'basis',htmlLabels:true,padding:14,nodeSpacing:36,rankSpacing:46}}}});
  const good=[];
  for(const el of nodes){{
   const src=el.textContent;
   try{{await mermaid.parse(src);good.push(el);}}
   catch(e){{el.classList.add('mmfail');el.innerHTML='<div class="mmfail-note">다이어그램을 표시할 수 없습니다.</div>';}}
  }}
  if(good.length){{
   await mermaid.run({{nodes:good}});
   good.forEach(el=>{{
    const svg=el.querySelector('svg');if(!svg)return;
    svg.style.maxWidth='100%';svg.style.width='100%';svg.removeAttribute('height');
    const box=el.closest('.fdiagram')||el.parentElement;
    if(!box||box.querySelector('.fdiagram-zoom'))return;
    box.classList.add('zoomable');
    const open=()=>{{const c=svg.cloneNode(true);
     c.style.width='auto';c.style.height='auto';c.style.maxWidth='94vw';c.style.maxHeight='88vh';
     window.__lb&&window.__lb.open(c);}};
    const btn=document.createElement('button');btn.type='button';btn.className='fdiagram-zoom';
    btn.textContent='🔍 확대';btn.addEventListener('click',ev=>{{ev.stopPropagation();open();}});
    box.appendChild(btn);
    el.addEventListener('click',open);
   }});
  }}
 }}catch(e){{
  nodes.forEach(m=>{{m.style.whiteSpace='pre-wrap';m.style.color='var(--ink-dim)';}});
 }}
 if(window.__fitAll) setTimeout(window.__fitAll,60);
}})();
</script>
<script>window.__NOTES={notes_json};window.__DECK={deck_json};</script>
<script>{PRESENT_JS}</script>
</body></html>"""


def render_deck_index_html(stem: str, versions: list[dict]) -> str:
    """한 PPT 폴더의 버전 히스토리 목록 (versions 는 최신순)."""
    rows = []
    for i, v in enumerate(versions):
        tag = " <span style='color:var(--accent)'>(최신)</span>" if i == 0 else ""
        rows.append(
            f'<li><a href="{v["id"]}.html">{escape(v["ts_human"])}</a>{tag}</li>'
        )
    return f"""<!DOCTYPE html>
<html lang="ko"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(stem)} — 버전 히스토리</title>
<style>{CSS}
body{{padding:0;}} .wrap{{max-width:760px;margin:0 auto;padding:48px 24px;}}
h1{{font-size:24px;margin:0 0 4px;}} .sub{{color:var(--muted);margin-bottom:28px;}}
ul.vers{{list-style:none;margin:0;padding:0;}}
ul.vers li{{padding:10px 0;border-bottom:1px solid var(--border);}}
ul.vers a{{color:var(--accent);text-decoration:none;font-weight:600;}}
.cta{{display:inline-block;margin:8px 0 18px;padding:10px 16px;background:var(--accent);
color:#fff;border-radius:10px;text-decoration:none;}}
</style></head><body>
<div class="wrap">
<h1>{escape(stem)}</h1>
<div class="sub">버전 히스토리 · 총 {len(versions)}개</div>
<a class="cta" href="latest.html">최신 버전 보기 &rarr;</a>
<ul class="vers">{''.join(rows)}</ul>
<p><a href="../index.html" style="color:var(--accent)">&larr; 전체 덱 목록</a></p>
</div></body></html>"""


def render_index_html(groups: dict, slugs: dict) -> str:
    """전체 덱 목록 (PPT 별 카드)."""
    cards = []
    for stem in sorted(groups, key=lambda s: groups[s][0]["ts_key"], reverse=True):
        vers = groups[stem]
        slug = slugs[stem]
        latest = vers[0]
        cards.append(
            f'<article class="card"><h2><a href="{slug}/latest.html">{escape(stem)}</a></h2>'
            f'<div class="meta">버전 {len(vers)}개 · 최신 {escape(latest["ts_human"])}</div>'
            f'<a href="{slug}/index.html" style="color:var(--accent);font-size:14px">버전 히스토리 &rarr;</a></article>'
        )
    return f"""<!DOCTYPE html>
<html lang="ko"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>스크롤 덱 목록</title>
<style>{CSS}
body{{padding:0;}}
.wrap{{max-width:880px;margin:0 auto;padding:48px 24px;}}
h1{{font-size:28px;margin:0 0 8px;}}
.sub{{color:var(--muted);margin-bottom:32px;}}
.card{{border:1px solid var(--border);border-radius:14px;padding:20px 24px;margin:16px 0;background:var(--card);}}
.card h2{{margin:0 0 4px;font-size:20px;}}
.card h2 a{{color:var(--fg);text-decoration:none;}}
.card .meta{{color:var(--muted);font-size:13px;margin-bottom:10px;}}
</style></head><body>
<div class="wrap">
<h1>스크롤 덱</h1>
<div class="sub">PPTX 내용을 HTML 로 변환한 스크롤형 발표 자료. PPT 별 폴더에서 버전 히스토리 관리.</div>
{''.join(cards) or '<p class="empty">아직 생성된 덱이 없습니다.</p>'}
</div></body></html>"""


# --------------------------------------------------------------------------- #
# 빌드
# --------------------------------------------------------------------------- #
def build(releases_dir: Path, out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    releases = [parse_release(f) for f in list_releases(releases_dir)]
    releases = [r for r in releases if r["pptx"]]

    groups: dict[str, list[dict]] = {}
    for r in releases:
        groups.setdefault(r["stem"], []).append(r)
    for stem in groups:
        groups[stem].sort(key=lambda r: r["ts_key"], reverse=True)

    slugs = {stem: slugify(stem) for stem in groups}

    count = 0
    for stem, vers in groups.items():
        folder = out_dir / slugs[stem]
        folder.mkdir(parents=True, exist_ok=True)
        for r in vers:
            deck = extract_deck(r["pptx"])
            titles = load_titles(r["titles_file"])
            content = load_content(r.get("content_file"))
            fragments = load_fragments(r.get("slides_dir"))
            theme = extract_theme(r["pptx"])
            html = render_version_html(r, deck, titles, vers, content, fragments, theme)
            (folder / f"{r['id']}.html").write_text(html, encoding="utf-8")
            if r is vers[0]:  # 최신
                (folder / "latest.html").write_text(html, encoding="utf-8")
            count += 1
            print(f"  + {slugs[stem]}/{r['id']}.html  ({len(deck)} slides)")
            # 프래그먼트가 참조하는 이미지(slides/img) 를 출력 폴더로 복사
            sd = r.get("slides_dir")
            if sd:
                img_dir = Path(sd) / "img"
                if img_dir.is_dir():
                    shutil.copytree(img_dir, folder / "img", dirs_exist_ok=True)
        (folder / "index.html").write_text(
            render_deck_index_html(stem, vers), encoding="utf-8")

    (out_dir / "index.html").write_text(render_index_html(groups, slugs), encoding="utf-8")
    print(f"빌드 완료: {count}개 버전, {len(groups)}개 덱 -> {out_dir}")


def main() -> None:
    ap = argparse.ArgumentParser(description="PPTX 내용을 스크롤형 HTML 덱으로 빌드")
    ap.add_argument("--releases", default="releases", help="릴리즈 폴더 (기본: releases)")
    ap.add_argument("--out", default="docs/decks", help="출력 폴더 (기본: docs/decks)")
    args = ap.parse_args()
    build(Path(args.releases), Path(args.out))


if __name__ == "__main__":
    main()
