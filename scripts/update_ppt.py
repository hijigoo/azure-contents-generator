"""기존 PPT 의 디자인(슬라이드 마스터·레이아웃·폰트)을 그대로 유지한 채
끝부분에 '최신 Azure 업데이트' 섹션(섹션 표지 + LLM 요약 + 항목별 슬라이드)을
추가/갱신한다.

LLM/외부 API 호출 없음. 인증 불필요.

흐름:
1. .cache/updates.json (fetch_azure_updates.py 결과) 로드
2. 노트에 마커가 있는 자동 생성 슬라이드 모두 제거 (멱등성)
3. 원본 PPT 의 슬라이드 레이아웃을 키워드로 탐색해 자동 매칭
   - 섹션 표지: Section / Header / Divider / Title
   - 본문      : Title and Content / Content
4. 다음 슬라이드를 순서대로 추가
   a. 섹션 커버: '🆕 최신 Azure 업데이트  ·  YYYY-MM-DD'
   b. (요약 있으면) 'TL;DR' 슬라이드 — LLM 요약 bullet
   c. items 상위 N개 — 한 슬라이드씩 (제목 + 발행일/출처 + 노트에 원문 링크)
   d. '데이터 소스' 마무리 슬라이드
5. zip 중복 entry 정리 (LibreOffice 호환)
"""
from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

AUTO_MARKER = "[auto-generated:azure-update]"


# ──────────────────────────────────────────────────────────────────────────
# 레이아웃 / 슬라이드 헬퍼
# ──────────────────────────────────────────────────────────────────────────

def pick_layout(prs: Presentation, keywords: list[str], fallback_idx: int = 1):
    """레이아웃 이름에 keywords 중 하나라도 포함되면 그것을 반환.
    못 찾으면 fallback_idx (없으면 마지막 layout)."""
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if any(k.lower() in name for k in keywords):
            return layout
    idx = min(fallback_idx, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def set_notes(slide, text: str) -> None:
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


def is_auto_slide(slide) -> bool:
    if not slide.has_notes_slide:
        return False
    try:
        return AUTO_MARKER in (slide.notes_slide.notes_text_frame.text or "")
    except Exception:
        return False


def remove_auto_slides(prs: Presentation) -> int:
    """노트에 AUTO_MARKER 가 있는 슬라이드를 모두 제거하고 개수를 반환."""
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    sld_ids = list(xml_slides)
    removed = 0
    for idx, slide in enumerate(list(prs.slides)):
        if is_auto_slide(slide):
            rId = sld_ids[idx].rId
            prs.part.drop_rel(rId)
            xml_slides.remove(sld_ids[idx])
            removed += 1
    return removed


def _strip_md(text: str) -> str:
    """마크다운 **bold** 마커만 제거한 평문."""
    return re.sub(r"\*\*([^*]+)\*\*", r"\1", text)


def _add_rich_paragraph(tf, text: str, *, first: bool, size_pt: int = 18) -> None:
    """**bold** 토큰을 인식해 run.bold = True 로 매핑."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    # 기존 run 비우기
    for r in list(p.runs):
        r.text = ""
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    first_part = True
    for part in parts:
        if not part:
            continue
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part
        run.font.size = Pt(size_pt)
        first_part = False
    if first_part:  # 모든 part 가 비어있던 경우
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_pt)


def _body_placeholder(slide):
    """제목이 아닌 첫 텍스트 placeholder 를 반환 (없으면 None)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            return ph
    return None


# ──────────────────────────────────────────────────────────────────────────
# 슬라이드 빌더
# ──────────────────────────────────────────────────────────────────────────

def add_section_cover(prs: Presentation, date_str: str) -> None:
    layout = pick_layout(
        prs,
        keywords=["section", "divider", "header", "title slide", "표지"],
        fallback_idx=0,
    )
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "🆕 최신 Azure 업데이트"
    body = _body_placeholder(slide)
    if body is not None:
        body.text_frame.text = f"{date_str} 기준 · 자동 생성"
        for run in body.text_frame.paragraphs[0].runs:
            run.font.size = Pt(20)
    set_notes(slide, f"{AUTO_MARKER}\n섹션 커버 (자동 생성, 다음 실행 시 교체됨).")


def add_summary_slide(prs: Presentation, summary_text: str) -> None:
    layout = pick_layout(
        prs,
        keywords=["title and content", "content", "본문", "title"],
        fallback_idx=1,
    )
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "TL;DR — 이번 주 핵심"

    bullets: list[str] = []
    for raw in summary_text.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith(("-", "*", "•")):
            line = line.lstrip("-*•").strip()
            bullets.append(line)
    if not bullets:
        bullets = [_strip_md(line) for line in summary_text.splitlines() if line.strip()]
    bullets = bullets[:6]

    body = _body_placeholder(slide)
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.0))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets or ["(요약 없음)"]):
        _add_rich_paragraph(tf, b, first=(i == 0), size_pt=18)

    set_notes(slide, f"{AUTO_MARKER}\nLLM 요약 (summary.md) 기반.")


def add_item_slide(prs: Presentation, item: dict) -> None:
    layout = pick_layout(
        prs,
        keywords=["title and content", "content", "본문", "title"],
        fallback_idx=1,
    )
    slide = prs.slides.add_slide(layout)

    title_txt = (item.get("title") or "Azure Update").strip()
    if len(title_txt) > 110:
        title_txt = title_txt[:107] + "…"
    if slide.shapes.title is not None:
        slide.shapes.title.text = title_txt

    published = (item.get("published") or "").strip()[:10]
    source = (item.get("source") or "").strip()
    link = (item.get("link") or "").strip()
    summary = (item.get("summary") or "").strip()
    if len(summary) > 260:
        summary = summary[:257] + "…"

    body = _body_placeholder(slide)
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.0))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    meta_parts = []
    if published:
        meta_parts.append(f"📅 {published}")
    if source:
        meta_parts.append(f"🏷 {source}")
    meta = "   ".join(meta_parts) or "📅 (date unknown)"
    _add_rich_paragraph(tf, meta, first=True, size_pt=14)

    if summary:
        _add_rich_paragraph(tf, summary, first=False, size_pt=18)

    if link:
        _add_rich_paragraph(tf, f"🔗 {link}", first=False, size_pt=12)

    notes = [AUTO_MARKER, ""]
    notes.append(f"제목: {item.get('title','')}")
    if published:
        notes.append(f"발행일: {published}")
    if source:
        notes.append(f"출처: {source}")
    if link:
        notes.append(f"원문: {link}")
    if summary:
        notes.append("")
        notes.append(summary)
    set_notes(slide, "\n".join(notes))


def add_sources_slide(prs: Presentation, items: list[dict]) -> None:
    layout = pick_layout(prs, keywords=["title and content", "content"], fallback_idx=1)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "📚 데이터 소스"
    body = _body_placeholder(slide)
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.0))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    seen = set()
    sources = []
    for it in items:
        s = (it.get("source") or "").strip()
        if s and s not in seen:
            seen.add(s)
            sources.append(s)
    lines = [f"• {s}" for s in sources] or ["• Azure / Microsoft 공식 RSS"]
    lines.append("")
    lines.append(f"⏱ 생성 시각: {datetime.utcnow().isoformat(timespec='seconds')}Z")
    for i, ln in enumerate(lines):
        _add_rich_paragraph(tf, ln, first=(i == 0), size_pt=14)
    set_notes(slide, f"{AUTO_MARKER}\n출처 목록 (자동 생성).")


def add_updates_section(
    prs: Presentation,
    items: list[dict],
    *,
    summary_text: str | None,
    limit: int,
) -> int:
    today = datetime.now().strftime("%Y-%m-%d")
    add_section_cover(prs, today)
    if summary_text:
        add_summary_slide(prs, summary_text)
    for item in items[:limit]:
        add_item_slide(prs, item)
    add_sources_slide(prs, items[:limit])
    # 섹션커버 + (요약 1) + items + 소스
    return 1 + (1 if summary_text else 0) + len(items[:limit]) + 1


# ──────────────────────────────────────────────────────────────────────────
# 기존 기능 (skill clean + dedupe)
# ──────────────────────────────────────────────────────────────────────────

def maybe_run_skill_clean(pptx: Path, skill_dir: Path) -> None:
    """skill 의 clean.py 가 있으면 호출해 임시 자산을 정리."""
    clean_script = skill_dir / "scripts" / "clean.py"
    if not clean_script.exists():
        return
    try:
        subprocess.run(
            [sys.executable, str(clean_script), str(pptx)],
            check=True,
            capture_output=True,
        )
        print(f"[update_ppt] skill clean.py 실행 완료")
    except subprocess.CalledProcessError as exc:
        # clean 단계는 실패해도 본 업데이트는 유지
        print(f"[update_ppt] skill clean.py 실패(무시): {exc.stderr.decode(errors='ignore')[:200]}")


def dedupe_pptx_zip(pptx_path: Path) -> int:
    """python-pptx 가 만든 .pptx 안에 같은 이름의 zip entry 가 여러 개 들어가는
    경우가 있다 (원본 .pptx 의 손상된 master.rels 등이 원인). LibreOffice/Keynote
    같은 엄격한 리더가 'source file could not be loaded' 로 거부하므로, 같은 이름은
    **마지막** entry 만 남기고 재압축한다. 제거된 entry 개수를 반환."""
    import shutil
    import zipfile

    with zipfile.ZipFile(pptx_path, "r") as zin:
        infos = zin.infolist()
        kept: dict[str, tuple[zipfile.ZipInfo, bytes]] = {}
        for info in infos:
            with zin.open(info) as fh:
                data = fh.read()
            kept[info.filename] = (info, data)
        removed = len(infos) - len(kept)

    if removed == 0:
        return 0

    tmp = pptx_path.with_suffix(pptx_path.suffix + ".dedup")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, (info, data) in kept.items():
            new_info = zipfile.ZipInfo(filename=name, date_time=info.date_time)
            new_info.compress_type = zipfile.ZIP_DEFLATED
            new_info.external_attr = info.external_attr
            zout.writestr(new_info, data)
    shutil.move(str(tmp), str(pptx_path))
    return removed


def main() -> None:
    parser = argparse.ArgumentParser(description="Azure 업데이트 슬라이드 기계 삽입")
    parser.add_argument("--input", required=True, help="원본 PPT 경로 (덮어씀)")
    parser.add_argument("--updates", required=True, help="updates.json 경로")
    parser.add_argument("--skill", required=True, help="anthropics/skills/skills/pptx 디렉터리")
    parser.add_argument("--limit", type=int, default=8, help="삽입할 항목 수")
    parser.add_argument("--summary", help="LLM(코파일럿) 생성 요약 마크다운 경로 (없으면 RSS bullet 사용)")
    parser.add_argument("--output", help="출력 경로 (기본: 원본 덮어쓰기)")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    pptx_path = Path(args.input)
    updates_path = Path(args.updates)
    skill_dir = Path(args.skill)
    output_path = Path(args.output) if args.output else pptx_path

    if not pptx_path.exists():
        sys.exit(f"입력 PPT를 찾을 수 없습니다: {pptx_path}")
    if not updates_path.exists():
        sys.exit(f"업데이트 JSON을 찾을 수 없습니다: {updates_path}")

    data = json.loads(updates_path.read_text(encoding="utf-8"))
    items = data.get("items", [])
    print(f"[update_ppt] 항목 {len(items)}건 중 상위 {args.limit}건 반영")

    if args.dry_run:
        for it in items[: args.limit]:
            print(f"  - {it.get('title', '')} ({it.get('published', '')[:10]})")
        return

    prs = Presentation(pptx_path)
    removed = remove_auto_slides(prs)
    if removed:
        print(f"[update_ppt] 기존 자동 생성 슬라이드 {removed}개 제거 (멱등)")

    summary_text = None
    if args.summary:
        sp = Path(args.summary)
        if sp.exists():
            summary_text = sp.read_text(encoding="utf-8")
            print(f"[update_ppt] LLM 요약 사용: {sp} ({len(summary_text)}자)")

    added = add_updates_section(prs, items, summary_text=summary_text, limit=args.limit)
    print(f"[update_ppt] 자동 생성 슬라이드 {added}장 추가 (섹션커버 + 요약 + 항목 + 소스)")

    prs.save(output_path)
    dedup = dedupe_pptx_zip(output_path)
    if dedup:
        print(f"[update_ppt] 중복 zip entry {dedup}개 정리됨 (LibreOffice 호환)")
    print(f"[update_ppt] 저장 완료: {output_path}")

    maybe_run_skill_clean(output_path, skill_dir)


if __name__ == "__main__":
    main()
