"""PPTX → PDF → 슬라이드별 PNG 변환.

GitHub에서 PPT 변경을 즉시 확인할 수 있도록 미리보기 자산을 생성한다.
- PDF: GitHub가 인라인 렌더링 지원
- PNG: PR 본문/README에 임베드해 슬라이드 단위 미리보기 제공

요구 사항(런타임): libreoffice, poppler-utils (pdftoppm)
"""
from __future__ import annotations

import argparse
import io
import os
import shutil
import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
DEFAULT_FONT = 18
SLIDE_SORT_FALLBACK = sys.maxsize


def require(cmd: str) -> str:
    path = shutil.which(cmd)
    if not path:
        sys.exit(f"필수 명령을 찾을 수 없습니다: {cmd}")
    return path


def slide_sort_key(path: Path) -> tuple[int, str]:
    try:
        return (int(path.stem.split("-")[-1]), path.name)
    except ValueError:
        return (SLIDE_SORT_FALLBACK, path.name)


def pptx_to_pdf(pptx: Path, out_dir: Path) -> Path:
    require("libreoffice")
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx),
        ],
        check=True,
    )
    pdf = out_dir / (pptx.stem + ".pdf")
    if not pdf.exists():
        sys.exit(f"PDF 변환 실패: {pdf}")
    return pdf


def pdf_to_png(pdf: Path, out_dir: Path, dpi: int = 110) -> list[Path]:
    require("pdftoppm")
    out_dir.mkdir(parents=True, exist_ok=True)
    prefix = out_dir / "slide"
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(prefix)],
        check=True,
    )
    return sorted(out_dir.glob("slide-*.png"), key=slide_sort_key)


def _relative_link(base_dir: Path, target: Path) -> str:
    if target.parent.resolve() == base_dir.resolve():
        return target.name
    return os.path.relpath(target.resolve(), start=base_dir.resolve()).replace("\\", "/")


def _font_candidates() -> list[Path]:
    candidates: list[Path] = []
    if os.environ.get("PPTX_PREVIEW_FONT"):
        candidates.append(Path(os.environ["PPTX_PREVIEW_FONT"]))
    candidates.extend(
        [
            Path.home() / ".local/share/fonts/NanumGothic.ttf",
            Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        ]
    )
    return candidates


def load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for path in _font_candidates():
        if path.exists():
            try:
                return ImageFont.truetype(str(path), max(size, 10))
            except OSError:
                continue
    return ImageFont.load_default()


def shape_box(shape, scale: float, offset_x: int = 0, offset_y: int = 0) -> tuple[int, int, int, int]:
    return (
        int(shape.left * scale) + offset_x,
        int(shape.top * scale) + offset_y,
        max(int(shape.width * scale), 1),
        max(int(shape.height * scale), 1),
    )


def wrap_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    max_width: int,
) -> list[str]:
    if not text:
        return [""]
    lines: list[str] = []
    for raw_line in text.splitlines() or [""]:
        current = ""
        for ch in raw_line:
            candidate = current + ch
            if current and draw.textlength(candidate, font=font) > max_width:
                lines.append(current.rstrip())
                current = "" if ch.isspace() else ch
            else:
                current = candidate
        lines.append(current.rstrip() or raw_line.strip() or "")
    return lines or [text]


def paragraph_text(paragraph) -> str:
    pieces = [run.text for run in paragraph.runs if run.text]
    if pieces:
        return "".join(pieces).strip()
    return paragraph.text.strip()


def paragraph_font_size(shape, paragraph) -> int:
    for run in paragraph.runs:
        if run.font.size:
            return max(int(run.font.size.pt), 10)
    if shape.is_placeholder and shape.placeholder_format.idx == 0:
        return 28
    return DEFAULT_FONT


def render_text_shape(canvas: Image.Image, draw: ImageDraw.ImageDraw, shape, scale: float, offset_x: int = 0, offset_y: int = 0) -> None:
    if not shape.has_text_frame:
        return
    x, y, width, height = shape_box(shape, scale, offset_x, offset_y)
    if width < 10 or height < 10:
        return

    inset = max(int(8 * scale), 6)
    cursor_y = y + inset
    max_width = max(width - inset * 2, 20)
    draw.rounded_rectangle((x, y, x + width, y + height), radius=max(int(6 * scale), 4), outline="#D0D7DE", width=1)

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph_text(paragraph)
        if not text:
            cursor_y += max(int(8 * scale), 8)
            continue
        prefix = ""
        if not (shape.is_placeholder and shape.placeholder_format.idx == 0):
            prefix = "• " if paragraph.level == 0 else "  " * paragraph.level + "◦ "
        size = paragraph_font_size(shape, paragraph)
        font = load_font(max(int(size * scale * 0.95), 12))
        for line in wrap_text(draw, prefix + text, font, max_width):
            if cursor_y >= y + height - inset:
                return
            draw.text((x + inset, cursor_y), line, fill="#111827", font=font)
            line_height = draw.textbbox((0, 0), line or "가", font=font)[3] + max(int(4 * scale), 4)
            cursor_y += line_height
        cursor_y += max(int(6 * scale), 4)


def render_picture_shape(canvas: Image.Image, shape, scale: float, offset_x: int = 0, offset_y: int = 0) -> None:
    x, y, width, height = shape_box(shape, scale, offset_x, offset_y)
    if width < 5 or height < 5:
        return
    try:
        picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
    except Exception:
        return
    contained = ImageOps.contain(picture, (width, height))
    paste_x = x + (width - contained.width) // 2
    paste_y = y + (height - contained.height) // 2
    canvas.paste(contained, (paste_x, paste_y))


def render_shape(canvas: Image.Image, draw: ImageDraw.ImageDraw, shape, scale: float, offset_x: int = 0, offset_y: int = 0) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_x, group_y, _, _ = shape_box(shape, scale, offset_x, offset_y)
        for child in shape.shapes:
            render_shape(canvas, draw, child, scale, group_x, group_y)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        render_picture_shape(canvas, shape, scale, offset_x, offset_y)
        return
    if shape.has_text_frame:
        render_text_shape(canvas, draw, shape, scale, offset_x, offset_y)


def render_locally(pptx: Path, out_dir: Path, dpi: int = 110) -> tuple[Path, list[Path]]:
    prs = Presentation(str(pptx))
    scale = dpi / EMU_PER_INCH
    width = max(int(prs.slide_width * scale), 1)
    height = max(int(prs.slide_height * scale), 1)
    pngs: list[Path] = []

    for idx, slide in enumerate(prs.slides, start=1):
        canvas = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(canvas)
        for shape in slide.shapes:
            render_shape(canvas, draw, shape, scale)
        png = out_dir / f"slide-{idx}.png"
        canvas.save(png)
        pngs.append(png)

    pdf = out_dir / f"{pptx.stem}.pdf"
    frames = [Image.open(path).convert("RGB") for path in pngs]
    try:
        if not frames:
            sys.exit("PNG 미리보기 생성 실패: 슬라이드가 없습니다.")
        frames[0].save(pdf, save_all=True, append_images=frames[1:])
    finally:
        for frame in frames:
            frame.close()
    return pdf, pngs


def main() -> None:
    parser = argparse.ArgumentParser(description="PPT 미리보기(PDF/PNG) 생성")
    parser.add_argument("--input", required=True, help="입력 PPTX 경로")
    parser.add_argument("--out", required=True, help="미리보기 출력 디렉터리")
    parser.add_argument("--dpi", type=int, default=110)
    args = parser.parse_args()

    pptx = Path(args.input)
    out_dir = Path(args.out)
    if not pptx.exists():
        sys.exit(f"입력 PPT를 찾을 수 없습니다: {pptx}")

    # 기존 미리보기 산출물만 제거(폴더 전체를 지우면 같은 폴더의 PPTX/요약이 함께 사라짐)
    out_dir.mkdir(parents=True, exist_ok=True)
    for stale in list(out_dir.glob("slide-*.png")) + list(out_dir.glob("*.pdf")):
        try:
            stale.unlink()
        except OSError:
            pass

    if shutil.which("libreoffice") and shutil.which("pdftoppm"):
        pdf = pptx_to_pdf(pptx, out_dir)
        pngs = pdf_to_png(pdf, out_dir, dpi=args.dpi)
    else:
        print("[render_ppt_previews] libreoffice/pdftoppm 미탑재 → python-pptx/Pillow 로컬 렌더링 폴백 사용")
        pdf, pngs = render_locally(pptx, out_dir, dpi=args.dpi)

    # README.md 생성: 미리보기 인덱스
    readme = out_dir / "README.md"
    pptx_link = _relative_link(out_dir, pptx)
    lines = [
        f"# 미리보기 — {pptx.name}",
        "",
        f"- 원본: [`{pptx_link}`](./{pptx_link})",
        f"- PDF: [`{pdf.name}`](./{pdf.name})",
        "",
        "## 슬라이드",
        "",
    ]
    for i, png in enumerate(pngs, start=1):
        lines.append(f"### Slide {i}")
        lines.append("")
        lines.append(f"![slide {i}](./{png.name})")
        lines.append("")
    readme.write_text("\n".join(lines), encoding="utf-8")

    print(f"[render_ppt_previews] PDF={pdf.name}, PNG={len(pngs)}장 생성 → {out_dir}")


if __name__ == "__main__":
    main()
