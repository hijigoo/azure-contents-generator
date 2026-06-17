#!/usr/bin/env python3
"""PPTX 슬라이드별 임베디드 사진을 추출해 최적화 WebP 로 저장하고 매니페스트를 출력.
배경(슬라이드 거의 전체를 덮는 그림)과 아주 작은 아이콘은 건너뛴다."""
import sys, hashlib
from pathlib import Path
from io import BytesIO
from pptx import Presentation
from pptx.util import Emu

EMU_PER_PX = 914400 / 96


def walk(shapes):
    for sh in shapes:
        st = sh.shape_type
        if st is not None and int(st) == 6:  # group
            yield from walk(sh.shapes)
        else:
            yield sh


def main(pptx_path, out_dir):
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    slide_area = sw * sh
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    from PIL import Image
    seen_bg = {}
    for i, slide in enumerate(prs.slides, 1):
        idx = 0
        for shp in walk(slide.shapes):
            try:
                if int(shp.shape_type) != 13:  # PICTURE
                    continue
            except Exception:
                continue
            try:
                blob = shp.image.blob
            except Exception:
                continue
            w = shp.width or 0
            h = shp.height or 0
            area_ratio = (w * h) / slide_area if slide_area else 0
            px_w = (w or 0) / EMU_PER_PX
            px_h = (h or 0) / EMU_PER_PX
            digest = hashlib.md5(blob).hexdigest()[:8]
            tag = ""
            if area_ratio > 0.82:
                tag = "BG-skip"
            elif max(px_w, px_h) < 90:
                tag = "tiny-skip"
            saved = ""
            if not tag:
                idx += 1
                try:
                    im = Image.open(BytesIO(blob))
                    im = im.convert("RGBA") if im.mode in ("RGBA","LA","P") else im.convert("RGB")
                    maxw = 1280
                    if im.width > maxw:
                        im = im.resize((maxw, round(im.height*maxw/im.width)))
                    buf = BytesIO()
                    im.save(buf, "WEBP", quality=86, method=5)
                    fn = f"s{i:02d}-{idx}.webp"
                    (out / fn).write_bytes(buf.getvalue())
                    saved = f"{fn} ({len(buf.getvalue())//1024}KB)"
                except Exception as e:
                    saved = f"ERR {e}"
            print(f"slide {i:2d} | {px_w:6.0f}x{px_h:4.0f}px | area {area_ratio:4.2f} | {digest} | {tag or 'KEEP'} {saved}")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
